import io
import json
import random
from copy import deepcopy
import re
from pptx import Presentation
from pptx.util import Inches
import pandas as pd
from utils import compare_with_multiple_gt, compare_with_multiple_gt_with_altid, get_category, gt_mention_category_count, read_jsonl
import wandb

class Evaluation:
    def __init__(self, eval_data,kb, candidate_key, multiple_gt_grag=None, error_analysis=False, for_retrieval=True):
        data = deepcopy(eval_data)
        self.data = data
        self.kb = kb
        self.candidate_key = candidate_key
        self.error_analysis = error_analysis
        self.for_retrieval = for_retrieval
        self.final_data = []
        self.mention_matched_at_k = {}
        self.cat_wise_gt_matched = {}
        self.mention_wise_triple_matched = {}   
        self.text_report = "\n\n\n"

        self.multiple_gt_grag_dict = {}
        if multiple_gt_grag:
            for item in multiple_gt_grag:
                 self.multiple_gt_grag_dict[item['sample_id']] = item

        for i in data:
            if 'retrieved_candidates' not in i:
                continue
            self.final_data.append(i)
        
        removed_count = len(data) - len(self.final_data) 
        if removed_count > 0:
            self.text_report+=f"Removed {removed_count} items since they did not have retrieved_candidates key\n"
        
        self.none_data = [] 
        self.not_none_data = []
        for item in self.data:
            ground_truth_id = item["ground_truth"]["id"]
            if ground_truth_id.lower() == "none":
                self.none_data.append(item) 
            else:
                self.not_none_data.append(item)
        self.text_report+=f"Number of items with ground truth 'None': {len(self.none_data)}\n"
        self.text_report+=f"Number of items with ground truth not 'None': {len(self.not_none_data)}\n"

    def get_report(self):
        self.mrr = self.calculate_mrr()
        self.recall_at_1 = self.calculate_recall_at_k_for_graph_retriever(1)
        # self.recall_at_5 = self.calculate_recall_at_k_for_graph_retriever(5)
        # self.recall_at_10 = self.calculate_recall_at_k_for_graph_retriever(10)
        self.recall_at_10 = self.calculate_recall_at_k_for_graph_retriever(64)
        # none_accuracy = self.get_none_accuracy()

        return self.not_none_data, self.none_data

                
    
    def get_none_accuracy(self):
        data = self.none_data
        correct_count = 0
        for item in data:
            retrieved_candidates = item["retrieved_candidates"]
            if len(retrieved_candidates) == 0:
                correct_count += 1
            else:
                try:
                    for i in retrieved_candidates:
                        if i["id"].lower() == "none" or not i["id"]:
                            correct_count += 1
                            break
                except Exception as e:
                    if isinstance(retrieved_candidates[0], str):
                        if retrieved_candidates[0].lower() == "none":
                            correct_count += 1
                        
        accuracy = correct_count / len(data)
        self.text_report+=f"Total {len(data)} items was actually NONE, among them {correct_count} items correctly predicted as NONE.\nSo, accuracy for items with ground truth 'None' : {accuracy}\n"
        return accuracy
    

    def compare_with_multiple_gt(self, candidate, gt_list):
        for gt_item in gt_list:
            if gt_item['id'] == candidate['id']:
                return True 
            
    def compare_with_multiple_gt_with_altid(self, candidate, gt_list):
        candidate_id = candidate['id']
        altids = self.kb[candidate_id]['altdiseaseid']
        for alid in altids:
            for gt_item in gt_list:
                if gt_item['id'] == alid:
                    return True 

    def calculate_mrr(self):
        data = self.not_none_data
        reciprocal_ranks = []
        for item in data:
            if self.multiple_gt_grag_dict:
                gt_list = self.multiple_gt_grag_dict[item['sample_id']]['ground_truth']
                rank = 0
                for idx, candidate in enumerate(item[self.candidate_key], start=1):
                    got_ranked = False
                    got_ranked = self.compare_with_multiple_gt(candidate, gt_list)
                    if got_ranked:
                        rank = idx
                        break
                    else:
                        got_ranked = self.compare_with_multiple_gt_with_altid(candidate, gt_list)
                        if got_ranked:
                            rank = idx
                            break

            else:
                ground_truth_id = item["ground_truth"]["id"]
                retrieved_candidates = item["retrieved_candidates"]
                rank = 0
                for idx, candidate in enumerate(retrieved_candidates, start=1):
                    if candidate["id"] == ground_truth_id:
                        rank = idx
                        break
            
            # Calculate reciprocal rank; if not found, reciprocal rank is 0
            reciprocal_rank = 1 / rank if rank > 0 else 0
            reciprocal_ranks.append(reciprocal_rank)
        
        # Calculate MRR by averaging all reciprocal ranks
        mrr = sum(reciprocal_ranks) / len(reciprocal_ranks)
        self.text_report+=f"\n\nCalculated MRR for {len(data)} items. \nSo, MRR : {mrr}\n\n"
        return mrr

    def calculate_recall_at_k(self, k):
        data = self.not_none_data
        relevant_count = 0
        for item in data:
            ground_truth_id = item["ground_truth"]["id"]
            retrieved_candidates = item["retrieved_candidates"][:k]  # Consider top-k candidates
            
            # Check if the relevant item is within the top-k candidates
            if any(candidate["id"] == ground_truth_id for candidate in retrieved_candidates):
                relevant_count += 1
        
        # Calculate recall@k
        recall_at_k = relevant_count / len(data)
        self.text_report+=f"\n\nFound {relevant_count} items out of {len(data)} within top-{k} candidates \nSo, recall@{k} : {relevant_count}/{len(data)}={recall_at_k}\n\n"
        return recall_at_k
    

    def category_specific_count(self, item, is_matched, k):
        mention_id=item["mention_id"]
        mention=item["mention"]
        title= item["ground_truth"]['title']
        candidates = ' | '.join([i['title'] for i in item["retrieved_candidates"]])
        
        mention_lower = mention.lower()
        title_lower = title.lower()
        
        if mention_lower == title_lower:
            self.cat_wise_gt_matched[k]["HO"]['count'] += 1
            self.cat_wise_gt_matched[k]["HO"]['items'].append({
                'mention_id':mention_id, 
                'mention':mention, 
                'gt_title':title, 
                'matched': is_matched,
                'retrieved_candidates': candidates
                })
            if is_matched:
                self.cat_wise_gt_matched[k]["HO"]['matched'] += 1
        elif mention_lower in title_lower and title_lower != mention_lower:
            self.cat_wise_gt_matched[k]["MINT"]['count'] += 1
            self.cat_wise_gt_matched[k]["MINT"]['items'].append({
                'mention_id':mention_id, 
                'mention':mention, 
                'gt_title':title, 
                'matched': is_matched,
                'retrieved_candidates': candidates
                })
            if is_matched:
                self.cat_wise_gt_matched[k]["MINT"]['matched'] += 1
        else:
            words_mention = set(mention_lower.split())
            words_title = set(title_lower.split())
            common_words = words_mention.intersection(words_title)
            if common_words:
                self.cat_wise_gt_matched[k]["LO"]['count'] += 1
                self.cat_wise_gt_matched[k]["LO"]['items'].append({
                'mention_id':mention_id, 
                'mention':mention, 
                'gt_title':title, 
                'matched': is_matched,
                'retrieved_candidates': candidates
                })
                if is_matched:
                    self.cat_wise_gt_matched[k]["LO"]['matched'] += 1
            else:
                self.cat_wise_gt_matched[k]["NO"]['count'] += 1
                self.cat_wise_gt_matched[k]["NO"]['items'].append({
                'mention_id':mention_id, 
                'mention':mention, 
                'gt_title':title, 
                'matched': is_matched,
                'retrieved_candidates': candidates
                })
                if is_matched:
                    self.cat_wise_gt_matched[k]["NO"]['matched'] += 1


    
    def calculate_recall_at_k_for_graph_retriever(self, k):
        data = self.not_none_data
        self.gt_matched = []
        self.gt_did_not_matched = []

        self.cat_wise_gt_matched[k] = {
            "HO":{'count':0, 'matched':0, 'items':[]}, 
            "MINT":{'count':0, 'matched':0, 'items':[]}, 
            "LO":{'count':0, 'matched':0, 'items':[]},
            "NO":{'count':0, 'matched':0, 'items':[]}
        }
        
        self.mention_matched_at_k[k] = []
        
        relevant_count = 0

        if k==10:
            self.mention_wise_triple_matched[k]=[]

        for item in data:
            ground_truth_id = item["ground_truth"]["id"]
            
            if not self.for_retrieval:
                unique_triple = {}
                for id in item["retrieved_candidates"]:
                    unique_triple[id['id']] = id
                item["unique_triple"] = unique_triple

            retrieved_candidates = list(item["unique_triple"].keys())
            top_k = retrieved_candidates[:k] # Consider top-k candidates
            
            is_matched = False
            # Check if the relevant item is within the top-k candidates
            for i, candidate in enumerate(top_k):
                # if item["mention"] == "liver disease":
                #     print(0)
                if ground_truth_id in candidate:
                    relevant_count += 1
                    is_matched = True
                    self.mention_matched_at_k[k].append(item["mention_id"])
                    
                    matched_dict = item["unique_triple"][candidate]
                    self.gt_matched.append({
                        'mention_id':item["mention_id"],
                        'mention':item["mention"],
                        'mention_context':item["mention_context"],
                        'ground_truth':item["ground_truth"],
                        'rank':i+1,
                        'matched_triple': {**{'triple':candidate}, **matched_dict},
                        'unique_triple':item["unique_triple"]
                        })
                    break

            if k==10:
                if is_matched:
                    triple_and_aug = []
                    for i in top_k:
                        if isinstance(item["unique_triple"], dict):
                            break

                        item["unique_triple"][i]['score']
                        triple_and_aug.append({'triple':i, 'aug':item["unique_triple"][i]})

                    self.mention_wise_triple_matched[k].append(
                        {
                            'mention_id':item["mention_id"],
                            'mention':item["mention"],
                            'ground_truth':item["ground_truth"],
                            f'top_{k}':triple_and_aug,
                        }
                    )
            
            # if not self.for_retrieval:
            self.category_specific_count(item,
                        is_matched=is_matched,
                        k=k
                    )
            
            if not is_matched:
                self.gt_did_not_matched.append({
                            'mention_id':item["mention_id"],
                            'mention':item["mention"],
                            'mention_context':item["mention_context"],
                            'ground_truth':item["ground_truth"],
                            'rank':-1,
                            'matched_triple':None,
                            'unique_triple':item["unique_triple"]
                            })
        
        # Calculate recall@k
        recall_at_k = relevant_count / len(data)
        self.text_report+=f"\n\nFound {relevant_count} items out of {len(data)} within top-{k} candidates \nSo, recall@{k} : {relevant_count}/{len(data)}={recall_at_k}\n\n"
        # self.text_report+=f"\n\n"
        for cat in self.cat_wise_gt_matched[k]:
            count = self.cat_wise_gt_matched[k][cat]['count']
            matched = self.cat_wise_gt_matched[k][cat]['matched']
            score = matched / count if count > 0 else 0
            self.text_report+=f"{cat}, count: {count}, matched: {matched}, score: {score}\n"

        return recall_at_k


class MultiGTEvaluation:
    def __init__(self, eval_data,kb, candidate_key, multiple_gt_grag=None, error_analysis=False, for_retrieval=True):
        data = deepcopy(eval_data)
        self.data = data
        self.kb = kb
        self.kb_id_list = []
        for entid in self.kb:
            self.kb_id_list.append(entid)
            ent = self.kb[entid]
            if 'altdiseaseid' in ent:
                self.kb_id_list+=ent['altdiseaseid']
            
        self.candidate_key = candidate_key
        self.error_analysis = error_analysis
        self.for_retrieval = for_retrieval
        self.final_data = []
        self.mention_matched_at_k = {}
        self.cat_wise_gt_matched = {}
        self.mention_wise_triple_matched = {}   
        self.text_report = "\n\n\n"

        self.multiple_gt_grag_dict = {}
        if multiple_gt_grag:
            for item in multiple_gt_grag:
                 self.multiple_gt_grag_dict[item['sample_id']] = item

        for i in data:
            if 'retrieved_candidates' not in i:
                continue
            self.final_data.append(i)
        
        removed_count = len(data) - len(self.final_data) 
        if removed_count > 0:
            self.text_report+=f"Removed {removed_count} items since they did not have retrieved_candidates key\n"
        
        self.none_data = [] 
        self.not_none_data = []
        for item in self.data:
            ground_truth_id = item["ground_truth"]
            if ground_truth_id:
                self.not_none_data.append(item)
            else:
                self.none_data.append(item)

        self.text_report+=f"Number of items with ground truth 'None': {len(self.none_data)}\n"
        self.text_report+=f"Number of items with ground truth not 'None': {len(self.not_none_data)}\n"

    def get_report(self):
        self.mrr = self.calculate_mrr()
        self.recall_at_1 = self.calculate_recall_at_k_for_graph_retriever(1)
        self.recall_at_5 = self.calculate_recall_at_k_for_graph_retriever(5)
        self.recall_at_32 = self.calculate_recall_at_k_for_graph_retriever(32)
        self.recall_at_64 = self.calculate_recall_at_k_for_graph_retriever(64)
        # none_accuracy = self.get_none_accuracy()

        return self.not_none_data, self.none_data

                
    
    def get_none_accuracy(self):
        data = self.none_data
        correct_count = 0
        for item in data:
            retrieved_candidates = item["retrieved_candidates"]
            if len(retrieved_candidates) == 0:
                correct_count += 1
            else:
                try:
                    for i in retrieved_candidates:
                        if i["id"].lower() == "none" or not i["id"]:
                            correct_count += 1
                            break
                except Exception as e:
                    if isinstance(retrieved_candidates[0], str):
                        if retrieved_candidates[0].lower() == "none":
                            correct_count += 1
                        
        accuracy = correct_count / len(data)
        self.text_report+=f"Total {len(data)} items was actually NONE, among them {correct_count} items correctly predicted as NONE.\nSo, accuracy for items with ground truth 'None' : {accuracy}\n"
        return accuracy
    

    def calculate_mrr(self):
        data = self.not_none_data
        reciprocal_ranks = []
        for item in data:
            if self.multiple_gt_grag_dict:
                gt_list = self.multiple_gt_grag_dict[item['sample_id']]['ground_truth']
                rank = 0
                for idx, candidate in enumerate(item[self.candidate_key], start=1):
                    got_ranked = False
                    got_ranked, gt_item = compare_with_multiple_gt(candidate, gt_list)
                    if got_ranked:
                        rank = idx
                        break
                    else:
                        got_ranked, gt_item, alt_id = compare_with_multiple_gt_with_altid(self.kb, candidate, gt_list)
                        if got_ranked:
                            rank = idx
                            break

            else:
                ground_truth_id = item["ground_truth"]["id"]
                retrieved_candidates = item["retrieved_candidates"]
                rank = 0
                for idx, candidate in enumerate(retrieved_candidates, start=1):
                    if candidate["id"] == ground_truth_id:
                        rank = idx
                        break
            
            # Calculate reciprocal rank; if not found, reciprocal rank is 0
            reciprocal_rank = 1 / rank if rank > 0 else 0
            reciprocal_ranks.append(reciprocal_rank)
        
        # Calculate MRR by averaging all reciprocal ranks
        mrr = sum(reciprocal_ranks) / len(reciprocal_ranks)
        self.text_report+=f"\n\nCalculated MRR for {len(data)} items. \nSo, MRR : {mrr}\n\n"
        return mrr

    def calculate_recall_at_k(self, k):
        data = self.not_none_data
        relevant_count = 0
        for item in data:
            ground_truth_id = item["ground_truth"]["id"]
            retrieved_candidates = item["retrieved_candidates"][:k]  # Consider top-k candidates
            
            # Check if the relevant item is within the top-k candidates
            if any(candidate["id"] == ground_truth_id for candidate in retrieved_candidates):
                relevant_count += 1
        
        # Calculate recall@k
        recall_at_k = relevant_count / len(data)
        self.text_report+=f"\n\nFound {relevant_count} items out of {len(data)} within top-{k} candidates \nSo, recall@{k} : {relevant_count}/{len(data)}={recall_at_k}\n\n"
        return recall_at_k
    

    def category_specific_count(self, item, matched_gt, gt_list, is_matched, k):
        sample_id=item["sample_id"]
        mention=item["mention"]
        candidates = ' | '.join([i['title'] for i in item[self.candidate_key]])

        for gti in  gt_list:
            if gti['id'] in self.kb_id_list:
                title = gti['title']
                break
            else:
                title = gti['title']




        cat = get_category(mention, title)
        self.cat_wise_gt_matched[k][cat]['count'] += 1
        self.cat_wise_gt_matched[k][cat]['items'].append({
                'mention_id':sample_id, 
                'mention':mention, 
                'gt_title':title, 
                'matched': is_matched,
                'retrieved_candidates': candidates
                })
        if is_matched:
                self.cat_wise_gt_matched[k][cat]['matched'] += 1

    def calculate_recall_at_k_for_graph_retriever(self, k):
        data = self.not_none_data
        self.gt_matched = []
        self.gt_did_not_matched = []

        self.cat_wise_gt_matched[k] = {
            "HO":{'count':0, 'matched':0, 'items':[]}, 
            "MINT":{'count':0, 'matched':0, 'items':[]}, 
            "LO":{'count':0, 'matched':0, 'items':[]},
            "NO":{'count':0, 'matched':0, 'items':[]}
        }
        
        self.mention_matched_at_k[k] = []
        
        matched_count = 0

        if k==10:
            self.mention_wise_triple_matched[k]=[]

        for item in data:
            gt_list = self.multiple_gt_grag_dict[item['sample_id']]['ground_truth']
            unique_candidates = {}
            for cand in item[self.candidate_key]:
                unique_candidates[cand['id']] = cand
            item["unique_candidates"] = unique_candidates

            retrieved_candidates = list(item["unique_candidates"].keys())
            top_k = retrieved_candidates[:k] # Consider top-k candidates
            
            is_matched = False
            matched_gt = {}
            # Check if the relevant item is within the top-k candidates
            for i, candidate in enumerate(top_k):
                candidate_info = item["unique_candidates"][candidate]
                got_match_for_multi_gt, matched_gt = compare_with_multiple_gt(candidate_info, gt_list)
                if got_match_for_multi_gt:
                    matched_count += 1
                    is_matched = True
                    self.mention_matched_at_k[k].append(item["sample_id"])
                    self.gt_matched.append({
                        'sample_id':item["sample_id"],
                        'mention':item["mention"],
                        'mention_context':item["mention_context"],
                        'ground_truth':gt_list,
                        'rank':i+1,
                        'matched_triple': {**{'triple':candidate}, **candidate_info},
                        'unique_candidates':item["unique_candidates"]
                        })
                    break
                else:
                    got_match_for_multi_altid, matched_gt, altid = compare_with_multiple_gt_with_altid(self.kb, candidate_info, gt_list)
                    if got_match_for_multi_altid:
                        matched_count += 1
                        is_matched = True
                        self.mention_matched_at_k[k].append(item["sample_id"])
                        self.gt_matched.append({
                            'sample_id':item["sample_id"],
                            'mention':item["mention"],
                            'mention_context':item["mention_context"],
                            'ground_truth':gt_list,
                            'rank':i+1,
                            'matched_triple': {**{'triple':candidate}, **candidate_info},
                            'unique_candidates':item["unique_candidates"]
                            })
                        break

            self.category_specific_count(item, matched_gt, gt_list, is_matched=is_matched, k=k)
            
            if not is_matched:
                self.gt_did_not_matched.append({
                            'sample_id':item["sample_id"],
                            'mention':item["mention"],
                            'mention_context':item["mention_context"],
                            'ground_truth':gt_list,
                            'rank':-1,
                            'matched_triple':None,
                            'unique_candidates':item["unique_candidates"]
                            })
        
        # Calculate recall@k
        recall_at_k = matched_count / len(data)
        self.text_report+=f"\n\nFound {matched_count} items out of {len(data)} within top-{k} candidates \nSo, recall@{k} : {matched_count}/{len(data)}={recall_at_k}\n\n"
        # self.text_report+=f"\n\n"
        for cat in self.cat_wise_gt_matched[k]:
            count = self.cat_wise_gt_matched[k][cat]['count']
            matched = self.cat_wise_gt_matched[k][cat]['matched']
            score = matched / count if count > 0 else 0
            self.text_report+=f"{cat}, count: {count}, matched: {matched}, score: {score}\n"

        return recall_at_k


def cat_eval(acc, filepath, kbpath):
    with open(filepath) as f:
        preds = json.load(f)
    with open(kbpath) as f:
        kb = json.load(f)
    converted = []
    
    for p in preds:
        mention_context = p['text'][0]
        gtid = p['mention_data']['kb_id'][0]
        gttitle = kb[gtid]['title']
        candidates = [{'id':i[0], 'title':''} for i in p['mention_data']['candidates'] if i[0]!=gtid ]
        retrieved_candidates = []
        if p['linked'] == 1.0:
            retrieved_candidates.append({'id':gtid, 'title':gttitle})
            retrieved_candidates+=random.sample(candidates, 9)
        else:
            retrieved_candidates+=random.sample(candidates, 10)

        mention = mention_context.split("[E1]")[1].split("[\E1]")[0].strip()
        
        d = {'mention_id' : '111111',
            'mention': mention,
            'mention_context':mention_context,
            'ground_truth': {'id':gtid, 'title':gttitle},
            'retrieved_candidates':retrieved_candidates}
        
        converted.append(d)

    
    evaluation = Evaluation(converted, for_retrieval=False)
    not_none_data, none_data = evaluation.get_report()

    report = f'Accuracy : {acc}\n\n{evaluation.text_report}'
    with open(f"{filepath.replace('.json', '_eval.txt')}", 'w') as f:
        f.write(report)
    with open(f"{filepath.replace('.json', '_category_info.json')}", 'w') as f:
        json.dump(evaluation.cat_wise_gt_matched, f, indent=1)

class PPTReport:
    def __init__(self):
        self.presentation = Presentation()

    def save_pptx(self, file_path):
        self.presentation.save(f"{file_path}.pptx")



    
    def add_new_slide(self, df, title):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[5])
        slide.shapes.title.text = title
        x, y, cx, cy = Inches(1), Inches(1.6), Inches(8), Inches(2)
        rows, cols = df.shape[0] + 1, df.shape[1]
        table = slide.shapes.add_table(rows, cols, x, y, cx, cy).table

        for table_col in table.columns:
            table_col.width = Inches(0.95)

        table.columns[0].width = Inches(8)

        for col_idx, column_name in enumerate(df.columns):
            table.cell(0, col_idx).text = column_name

        for row_idx, row in df.iterrows():
            for col_idx, value in enumerate(row):
                table.cell(row_idx + 1, col_idx).text = str(value)


        
class ReportMaker():
    def __init__(self, onto, sptitname, 
                 
                bi_report_save_to, bi_report_start_end,
                bi_firt_row_text, 

                cross_report_save_to,cross_report_start_end,
                cross_firt_row_text, 

                cross_before_ft_dir=None,
                bi_before_ft_dir=None
                
        ): 

        self.sptitname = sptitname
        self.onto = onto

        self.bi_report_save_to = bi_report_save_to
        self.bi_report_start_end = bi_report_start_end
        self.bi_firt_row_text = bi_firt_row_text
        self.bi_before_ft_dir = bi_before_ft_dir
        
        self.cross_report_start_end = cross_report_start_end
        self.cross_report_save_to = cross_report_save_to
        self.cross_firt_row_text = cross_firt_row_text
        self.cross_before_ft_dir = cross_before_ft_dir


    def get_cat_wise_score(self, report, row):
        lines = report.split("\n")
        for line in lines:
            if "HO," in line:
                l = line.split('score:')
                lr = l[1].strip()
                row['HO'] =  round(float(lr) * 100, 2)
                ho_count = int(re.search(r'count:\s*(\d+)', line).group(1))
                ho_matched = int(re.search(r'matched:\s*(\d+)', line).group(1))
            elif "MINT," in line:
                l = line.split('score:')
                lr = l[1].strip()
                row['MINT'] = round(float(lr) * 100, 2)
                mint_count = int(re.search(r'count:\s*(\d+)', line).group(1))
                mint_matched = int(re.search(r'matched:\s*(\d+)', line).group(1))
            elif "LO," in line:
                l = line.split('score:')
                lr = l[1].strip()
                row['LO'] = round(float(lr) * 100, 2)
            elif "NO," in line:
                l = line.split('score:')
                lr = l[1].strip()
                row['NO'] = round(float(lr) * 100, 2)

        # merge HO & MINT
        ho_mint_count = ho_count+mint_count
        ho_mint_matched = ho_matched+mint_matched

        row['HO'] = round((ho_mint_matched/ho_mint_count) * 100, 2)
        del row['MINT']




    def save_to_wandb(self, log_data):
        for entry in log_data:
            wandb.log({
                "Overall": entry["Overall"],
                "HO": entry["HO"],
                # "MINT": entry["MINT"],
                "LO": entry["LO"],
                "NO": entry["NO"],
                'MRR':entry["MRR"],
            }, step=entry["Epoch"])

    def get_mrr_and_overall_for_best_model(self, epoch, report, best_epoch, reacall_at_i):
        mrr = None
        was_best = False
        lines = report.split("\n")
        for line in lines:
            if "MRR :" in line:
                l = line.split('MRR :')
                lr = l[1].strip()
                mrr =  round(float(lr) * 100, 2)
            
            if reacall_at_i in line:
                l = line.split(reacall_at_i)
                lr = l[1].split('=')
                overall_acc =  round(float(lr[1]) * 100, 2)
                if overall_acc>best_epoch['acc']:
                    best_epoch = {'acc':overall_acc, 'report':report, 'epoch':epoch, 'MRR':mrr}
                    was_best = True
        return mrr, overall_acc, best_epoch, was_best
         
    def get_overall(self, report, reacall_at_i):
        lines = report.split("\n")
        for line in lines:
            if reacall_at_i in line:
                l = line.split(reacall_at_i)
                lr = l[1].split('=')
                overall_acc =  round(float(lr[1]) * 100, 2)
                return overall_acc
            

    def get_best_epoch_and_send_all_wandb(self, exp, eval_text_name, text_eval_file_dir, epoch_start, epoch_end):
        report_all = []
        if self.cross_before_ft_dir:
            row = self.before_ft_results
            report_all.append(row)
        main_recall = 'recall@1'
        re_st_end = self.cross_report_start_end[main_recall]
        best_epoch = {'acc':0, 'report':'', 'epoch':None, 'MRR':None}
        entire_report_best_model = ''
        report_not_found = []
        for i in range(epoch_start, epoch_end+1):
            epoch_dir = f'{text_eval_file_dir}{exp}/epoch_{i}/'
            
            try:
                with open(f"{epoch_dir}{eval_text_name}", 'r') as f:
                    entire_report = f.read()
                    report = entire_report.split(re_st_end[0])[1]
                    report = report.split(re_st_end[1])[0]
            except FileNotFoundError:
                print(f"No report found : \n{epoch_dir}")
                report_not_found.append(i)

            if epoch_start==0:
                epoch = i+1
            elif epoch_start==1:
                epoch = i

            mrr, overall_acc, best_epoch, was_best = self.get_mrr_and_overall_for_best_model(epoch, report, best_epoch, 
                                                                                   main_recall)
            if was_best:
                entire_report_best_model = entire_report
                                                                                   

            row = {'Epoch':epoch,'Overall': overall_acc, 'MRR':mrr}
            self.get_cat_wise_score(report, row)
            report_all.append(row)

        # wandb.init(dir="wandb_logs", 
        #            project=f"FT-reranker-BLINK-{self.onto}", 
        #            name=exp+'20neg')
        

        # self.save_to_wandb(report_all)
        # wandb.finish()

        self.get_other_reacalls(entire_report_best_model, best_epoch)

        return best_epoch, report_not_found
    
    
    



    def get_last_epoch_and_send_all_wandb(self, exp, eval_text_name, text_eval_file_dir, epoch_start, epoch_end):
        report_all = []
        if self.cross_before_ft_dir:
            row = self.before_ft_results
            report_all.append(row)

        last_epoch = 0
        epoch_count = 0
        for i in range(epoch_start, epoch_end+1):
            epoch_dir = f'{text_eval_file_dir}{exp}/epoch_{i}/'
            try:
                with open(f"{epoch_dir}{eval_text_name}", 'r') as f:
                    entire_report = f.read()
                    epoch_count+=1
            except FileNotFoundError:
                last_epoch = i-1
                break

        last_epoch = 2 # fixed epoch
        print(f'Last epoch fixed to : {last_epoch+1} : {exp}')

        main_recall = 'recall@1'
        re_st_end = self.cross_report_start_end[main_recall]
        best_epoch = {'acc':0, 'report':'', 'epoch':None, 'MRR':None}
        entire_report_best_model = ''
        report_not_found = []


        
        epoch_start, epoch_end = last_epoch, last_epoch
        for i in range(epoch_start, epoch_end+1):
            epoch_dir = f'{text_eval_file_dir}{exp}/epoch_{i}/'
            
            try:
                with open(f"{epoch_dir}{eval_text_name}", 'r') as f:
                    entire_report = f.read()
                    report = entire_report.split(re_st_end[0])[1]
                    report = report.split(re_st_end[1])[0]
            except FileNotFoundError:
                print(f"No report found : \n{epoch_dir}")
                report_not_found.append(i)

            if epoch_start==0:
                epoch = i+1
            elif epoch_start==1:
                epoch = i
            else:
                epoch = i


            mrr, overall_acc, best_epoch, was_best = self.get_mrr_and_overall_for_best_model(epoch, report, best_epoch, 
                                                                                   main_recall)
            if was_best:
                entire_report_best_model = entire_report
                                                                                   

            row = {'Epoch':epoch,'Overall': overall_acc, 'MRR':mrr}
            self.get_cat_wise_score(report, row)
            report_all.append(row)

        # wandb.init(dir="wandb_logs", 
        #            project=f"FT-reranker-BLINK-{self.onto}", 
        #            name=exp+'20neg')
        

        # self.save_to_wandb(report_all)
        # wandb.finish()

        self.get_other_reacalls(entire_report_best_model, best_epoch)

        return best_epoch, report_not_found, epoch_count
    
    

    def get_specific_epoch_for_specifc_exp(self, exp, eval_text_name, text_eval_file_dir, epoch_start, epoch_end):
        report_all = []
        if self.cross_before_ft_dir:
            row = self.before_ft_results
            report_all.append(row)




        main_recall = 'recall@1'
        re_st_end = self.cross_report_start_end[main_recall]
        best_epoch = {'acc':0, 'report':'', 'epoch':None, 'MRR':None}
        entire_report_best_model = ''
        report_not_found = []
        exp_epoch_dict = {
            "(m1_e1)U(m3_e1)_multi_primeU(m4_e2)_multi_prime":0,
            "(m1_e1)U(m3_e1)_multi_prime_rm_sm_eU(m4_e2)_multi_prime_rm_sm_e":1,
            "synonym":3,
            "synonymU(m1_e1)U(m3_e1)_multi_prime_rm_sm_eU(m4_e2)_multi_prime_rm_sm_e":3
        }
        epoch_start, epoch_end = exp_epoch_dict[exp], exp_epoch_dict[exp]
        print(f'Epoch : {exp_epoch_dict[exp]+1} : {exp}')
        for i in range(epoch_start, epoch_end+1):
            epoch_dir = f'{text_eval_file_dir}{exp}/epoch_{i}/'
            
            try:
                with open(f"{epoch_dir}{eval_text_name}", 'r') as f:
                    entire_report = f.read()
                    report = entire_report.split(re_st_end[0])[1]
                    report = report.split(re_st_end[1])[0]
            except FileNotFoundError:
                print(f"No report found : \n{epoch_dir}")
                report_not_found.append(i)

            epoch = i


            mrr, overall_acc, best_epoch, was_best = self.get_mrr_and_overall_for_best_model(epoch, report, best_epoch, 
                                                                                   main_recall)
            if was_best:
                entire_report_best_model = entire_report
                                                                                   

            row = {'Epoch':epoch,'Overall': overall_acc, 'MRR':mrr}
            self.get_cat_wise_score(report, row)
            report_all.append(row)

        # wandb.init(dir="wandb_logs", 
        #            project=f"FT-reranker-BLINK-{self.onto}", 
        #            name=exp+'20neg')
        

        # self.save_to_wandb(report_all)
        # wandb.finish()

        self.get_other_reacalls(entire_report_best_model, best_epoch)

        return best_epoch, report_not_found, exp_epoch_dict[exp]
    
    


    
    
    def get_other_reacalls(self, entire_report_best_model, best_epoch):
        recalls = {}
        for reacall_at_i in self.cross_report_start_end:
            if reacall_at_i == 'recall@1':
                continue
            report_start_end_text = self.cross_report_start_end[reacall_at_i]
            recall_report = entire_report_best_model.split(report_start_end_text[0])[1]
            recall_report = recall_report.split(report_start_end_text[1])[1]
            recall_report = recall_report.split(report_start_end_text[2])[0]
            recall_row = {'Overall': self.get_overall(recall_report, reacall_at_i)}
            self.get_cat_wise_score(recall_report, recall_row)
            recalls[reacall_at_i] = recall_row
        best_epoch['recalls'] = recalls
        
    def bi_report_before_fine_tune_result(self, eval_text_name):

        main_recall = 'recall@1'
        re_st_end = self.bi_report_start_end[main_recall]
        try:
            with open(f"{self.bi_before_ft_dir}/{eval_text_name}", 'r') as f:
                entire_report = f.read()
                report = entire_report.split(re_st_end[0])[1]
                report = report.split(re_st_end[1])[0]
        except FileNotFoundError:
            print(f"No report found : \n{self.bi_before_ft_dir}{eval_text_name}")
            
        mrr, overall_acc, best_epoch, was_best = self.get_mrr_and_overall_for_best_model(0, report,  
        {'acc':0, 'report':'', 'epoch':None, 'MRR':None}, main_recall)
        row = {'Epoch':0,'Overall': overall_acc, 'MRR':mrr}
        self.get_cat_wise_score(report, row)

        recalls = {}
        for reacall_at_i in self.bi_report_start_end:
            if reacall_at_i == 'recall@1':
                continue
            report_start_end_text = self.bi_report_start_end[reacall_at_i]
            if reacall_at_i == 'recall@64':
                recall_report = entire_report.split(report_start_end_text[0])[1]
            else:
                recall_report = entire_report.split(report_start_end_text[0])[1]
                recall_report = recall_report.split(report_start_end_text[1])[1]
                recall_report = recall_report.split(report_start_end_text[2])[0]

            recall_row = {'Overall': self.get_overall(recall_report, reacall_at_i)}
            self.get_cat_wise_score(recall_report, recall_row)
            recalls[reacall_at_i] = recall_row


        row['recalls'] = recalls
        
        return row

    def bi_fine_tune_result(self, text_eval_file_dir, exp, eval_text_name, epoch_start, epoch_end):

        main_recall = 'recall@1'
        re_st_end = self.bi_report_start_end[main_recall]

        all_epochs = []
        best_recall_64 = 0
        
        for i in range(epoch_start, epoch_end+1):
            try:
                report_file = f"{text_eval_file_dir}{exp}/epoch_{i}/top64_candidates/{eval_text_name}"
                with open(report_file, 'r') as f:
                    entire_report = f.read()
                    report = entire_report.split(re_st_end[0])[1]
                    report = report.split(re_st_end[1])[0]
            except FileNotFoundError:
                print(f"No report found : \n{report_file}")
                
            mrr, overall_acc, _, was_best = self.get_mrr_and_overall_for_best_model(0, report,  
            {'acc':0, 'report':'', 'epoch':i, 'MRR':None}, main_recall)
            
            if epoch_start==0:
                current_epoch = i+1
            else:
                current_epoch = i

            row = {'Epoch':i,'Overall': overall_acc, 'MRR':mrr}
            self.get_cat_wise_score(report, row)

            recalls = {}
            for reacall_at_i in self.bi_report_start_end:
                if reacall_at_i == 'recall@1':
                    continue
                report_start_end_text = self.bi_report_start_end[reacall_at_i]
                if reacall_at_i == 'recall@64':
                    recall_report = entire_report.split(report_start_end_text[0])[1]
                else:
                    recall_report = entire_report.split(report_start_end_text[0])[1]
                    recall_report = recall_report.split(report_start_end_text[1])[1]
                    recall_report = recall_report.split(report_start_end_text[2])[0]
                recall_row = {'Overall': self.get_overall(recall_report, reacall_at_i)}
                self.get_cat_wise_score(recall_report, recall_row)
                recalls[reacall_at_i] = recall_row

            row['recalls'] = recalls
            all_epochs.append(row)

            recall_64 = row['recalls']['recall@64']['Overall']
            if recall_64 > best_recall_64:
                best_recall_64=recall_64
                best_epoch = row

        
        return best_epoch

    def cross_report_before_fine_tune_result(self, eval_text_name):

        main_recall = 'recall@1'
        re_st_end = self.cross_report_start_end[main_recall]
        try:
            with open(f"{self.cross_before_ft_dir}{eval_text_name}", 'r') as f:
                entire_report = f.read()
                report = entire_report.split(re_st_end[0])[1]
                report = report.split(re_st_end[1])[0]
        except FileNotFoundError:
            print(f"No report found : \n{self.cross_before_ft_dir}")
            
        mrr, overall_acc, best_epoch, was_best = self.get_mrr_and_overall_for_best_model(0, report,  
        {'acc':0, 'report':'', 'epoch':None, 'MRR':None}, main_recall)
        row = {'Epoch':0,'Overall': overall_acc, 'MRR':mrr}
        self.get_cat_wise_score(report, row)

        recalls = {}
        for reacall_at_i in self.cross_report_start_end:
            if reacall_at_i == 'recall@1':
                continue
            report_start_end_text = self.cross_report_start_end[reacall_at_i]

            recall_report = entire_report.split(report_start_end_text[0])[1]
            recall_report = recall_report.split(report_start_end_text[1])[1]
            recall_report = recall_report.split(report_start_end_text[2])[0]
            recall_row = {'Overall': self.get_overall(recall_report, reacall_at_i)}
            self.get_cat_wise_score(recall_report, recall_row)
            recalls[reacall_at_i] = recall_row

        row['recalls'] = recalls
        
        return row
    
    def get_cat_stats(self, exp):

        with open(f"data/{self.onto}/blink_format/{self.sptitname}/{exp}/train_category_count.json") as f:
            counts = json.load(f)
        ho = counts['HO']['count']+counts['MINT']['count']
        stats = f"H:{ho},L:{counts['LO']['count']},N:{counts['NO']['count']},UM:{counts['unique_mention_count']},UGT:{counts['unique_gt_count']}"
        return stats
    
        
    def get_pos_stats(self, exp, bienc_model_dir, epoch):
        with open(f'models/{self.onto}/biencoder/{self.sptitname}/{bienc_model_dir}/{exp}/diagnostics.json') as f:
            digs = json.load(f)
            for d in digs:
                if d['epoch'] == epoch:
                    pos = d['avg_pos_per_sample']
                    neg = d['avg_neg_per_sample']
                    steps = d['num_steps_in_this_epoch']
                    warmup = d['total_warmup_steps']
                    bs = d['train_batch_size']
                    return pos, neg, steps, warmup, bs


    # def make_report_for_crossencoder(self,text_eval_file_dir, eval_text_name, exps, epoch_start, epoch_end,
    #                                    exp_from=None):
    #     ppt_report = PPTReport()
    #     recall_at_1_reports = []
    #     recall_at_n_reports = {}
    #     for reacall_at_i in self.cross_report_start_end:
    #         recall_at_n_reports[reacall_at_i]=[]

    #     if self.cross_before_ft_dir:
    #         self.before_ft_results = self.cross_report_before_fine_tune_result(eval_text_name)
    #         row = {**{f'Settings : {self.cross_firt_row_text}':f" Original model (without fine-tune)​"},
    #             **self.before_ft_results }
    #         del row['Epoch']
            

    #         for recall_i in row['recalls']:
    #             recall_i_data = {"Settings : ": "Original model (without fine-tune)",
    #                          "Epoch": None}
    #             for key in row['recalls'][recall_i]:
    #                 recall_i_data[key] = row['recalls'][recall_i][key]
    #             recall_at_n_reports[recall_i].append(recall_i_data)


    #         del row['recalls']
    #         recall_at_1_reports.append(row)

    #     self.best_epoch_dir = {}
    #     for exp in exps:
    #         best_epoch = self.get_best_epoch_and_send_all_wandb(exp, eval_text_name, 
    #                         text_eval_file_dir, epoch_start, epoch_end)
            

    #         if epoch_start==0:
    #             best_epoch_num = best_epoch['epoch']-1
    #         elif epoch_start==1:
    #             best_epoch_num = best_epoch['epoch']

    #         self.best_epoch_dir[exp] = f"epoch_{best_epoch_num}"
            
    #         # if not exp_from:
    #         exp_from = exp

    #         num_samples =len( read_jsonl(f'data/{self.onto}/blink_format/{self.sptitname}/{exp_from}/train.jsonl'))
                
    #         cat_stats = self.get_cat_stats(exp_from)

    #         setting_name = f'Settings : {self.cross_firt_row_text}'
    #         row = {setting_name:f"{exp}, ({num_samples})({cat_stats}), epoch {best_epoch['epoch']}", 
    #             'Overall':best_epoch['acc'], 'MRR':best_epoch['MRR']}
    #         self.get_cat_wise_score(best_epoch['report'], row)
            
            

    #         for recall_i in best_epoch['recalls']:
    #             recall_i_data = {"Settings : ":row[setting_name]}
    #             for key in best_epoch['recalls'][recall_i]:
    #                 recall_i_data[key] = best_epoch['recalls'][recall_i][key]
                
    #             recall_at_n_reports[recall_i].append(recall_i_data)

            
    #         recall_at_1_reports.append(row)


        
    #     df = pd.DataFrame(recall_at_1_reports)
    #     title = f"{self.onto.upper()}:  recall@1"
    #     ppt_report.add_new_slide(df, title)
    
    #     for reacall_at_i in self.cross_report_start_end:
    #         if reacall_at_i == 'recall@1':
    #             continue

    #         df = pd.DataFrame(recall_at_n_reports[reacall_at_i])
    #         title = f"{self.onto.upper()}:  {reacall_at_i}"
    #         ppt_report.add_new_slide(df, title)


    #     ppt_file_path = f"{self.cross_report_save_to}/evaluation_report"
    #     ppt_report.save_pptx(ppt_file_path)


     
    #     json_path = f"{self.cross_report_save_to}/evaluation_report.json"
    #     with open(json_path, 'w') as f:
    #         json.dump(recall_at_n_reports, f, indent=2)

    def make_report_for_crossencoder(self, text_eval_file_dir, eval_text_name, exps,
                                    epoch_start, epoch_end, exp_from=None):
        """Build PPTX (one slide per recall@k) and a matching JSON dump.

        Note: recall@k>1 cat-wise values (HO/LO/NO) come from
        self.get_best_epoch_and_send_all_wandb(...). If those look wrong (e.g.
        pinned at 100.0), the bug is in that upstream call, not here — fix it
        where best_epoch['recalls'][recall@k] is populated.
        """
        ppt_report = PPTReport()
        recall_at_1_reports = []
        recall_at_n_reports = {k: [] for k in self.cross_report_start_end}

        settings_key = f'Settings : {self.cross_firt_row_text}'

        # ---- No-fine-tune row ----
        if self.cross_before_ft_dir:
            self.before_ft_results = self.cross_report_before_fine_tune_result(eval_text_name)
            row = {**{settings_key: " Original model (without fine-tune)"},
                **self.before_ft_results}
            # No-fine-tune has no training epoch; render as None (NaN in PPTX).
            row['Epoch'] = None

            for recall_i in row['recalls']:
                recall_i_data = {"Settings : ": "Original model (without fine-tune)",
                                "Epoch": None}
                for key in row['recalls'][recall_i]:
                    recall_i_data[key] = row['recalls'][recall_i][key]
                recall_at_n_reports[recall_i].append(recall_i_data)

            del row['recalls']
            recall_at_1_reports.append(row)

        report_not_found = {}
        # ---- Per-experiment rows ----
        self.best_epoch_dir = {}
        for exp in exps:
            best_epoch, report_nf, epoch_count = self.get_last_epoch_and_send_all_wandb(
                exp, eval_text_name, text_eval_file_dir, epoch_start, epoch_end)

            # best_epoch, report_nf, epoch_count = self.get_specific_epoch_for_specifc_exp(
            #     exp, eval_text_name, text_eval_file_dir, epoch_start, epoch_end)
            
            report_not_found[exp] = report_nf
            
            best_epoch['epoch'] = best_epoch['epoch']+1

            exp_from = exp
            num_samples = len(read_jsonl(
                f'data/{self.onto}/blink_format/{self.sptitname}/{exp_from}/train.jsonl'))
            cat_stats = self.get_cat_stats(exp_from)

            long_setting = f"{exp}, ({num_samples})({cat_stats})"

            row = {settings_key: long_setting,
                'Epoch':   best_epoch['epoch'],
                'Overall': best_epoch['acc'],
                'MRR':     best_epoch['MRR']}
            self.get_cat_wise_score(best_epoch['report'], row)   # adds HO/LO/NO



            for recall_i in best_epoch['recalls']:
                recall_i_data = {"Settings : ": long_setting,
                                "Epoch":best_epoch['epoch']}
                for key in best_epoch['recalls'][recall_i]:
                    recall_i_data[key] = best_epoch['recalls'][recall_i][key]
                recall_at_n_reports[recall_i].append(recall_i_data)

            recall_at_1_reports.append(row)

        will_print = False
        for rnf in report_not_found:
            if report_not_found[rnf]:
                will_print = True

        if will_print:
            print('report_not_found')
            print(json.dumps(report_not_found))

        # ---- PPTX (one slide per recall@k) ----
        df = pd.DataFrame(recall_at_1_reports)
        ppt_report.add_new_slide(df, f"{self.onto.upper()}:  recall@1")

        for recall_i in self.cross_report_start_end:
            if recall_i == 'recall@1':
                continue
            df = pd.DataFrame(recall_at_n_reports[recall_i])
            # ppt_report.add_new_slide(df, f"{self.onto.upper()}:  {recall_i}")

        ppt_report.save_pptx(f"{self.cross_report_save_to}/evaluation_report")
        del recall_at_n_reports["recall@5"]
        # ---- JSON dump ----
        # Backfill recall@1: the per-recall loop above iterates row['recalls'],
        # which only contains k>1. recall@1 metrics live at the top level of
        # each row, so we project them into the same shape as recall@k entries.
        recall_at_n_reports['recall@1'] = [
            {
                "Settings : ": r[settings_key],
                "Epoch":   r.get("Epoch"),
                "Overall": r.get("Overall"),
                "MRR":     r.get("MRR"),
                "HO":      r.get("HO"),
                "LO":      r.get("LO"),
                "NO":      r.get("NO"),
            }
            for r in recall_at_1_reports
        ]

        
        json_path = f"{self.cross_report_save_to}/evaluation_report.json"
        with open(json_path, 'w') as f:
            json.dump(recall_at_n_reports, f, indent=2)




    def make_report_for_biencoder(self,text_eval_file_dir, eval_text_name, exps, epoch_start, epoch_end,
                                       exp_from=None, bienc_model_dir=None):
        
        
        ppt_report = PPTReport()
        recall_at_1_reports = []
        recall_at_n_reports = {}
        for reacall_at_i in self.bi_report_start_end:
            recall_at_n_reports[reacall_at_i]=[]

        if self.bi_before_ft_dir:
            self.bi_before_ft_results = self.bi_report_before_fine_tune_result(eval_text_name)
            row = {**{f'Settings : {self.bi_firt_row_text}':f" Original model (without fine-tune)​"},
                **self.bi_before_ft_results }
            del row['Epoch']
            
            for recall_i in row['recalls']:
                recall_i_data = {"Settings : ":"Original model (without fine-tune)​"}
                for key in row['recalls'][recall_i]:
                    recall_i_data[key] = row['recalls'][recall_i][key]
                recall_at_n_reports[recall_i].append(recall_i_data)


            del row['recalls']
            recall_at_1_reports.append(row)

        self.best_epoch_dir = {}

        # ft_results = self.bi_fine_tune_result(eval_text_name)
        # ft_row = {**{f'Settings : {self.bi_firt_row_text}':f" Original model (without fine-tune)​"},
        #     **ft_results }
        # del ft_row['Epoch']
        
        for exp in exps:
            
            best_epoch = self.bi_fine_tune_result(text_eval_file_dir, exp, eval_text_name, epoch_start, epoch_end)
            exp_from = exp
            num_samples = len(read_jsonl(f'data/{self.onto}/blink_format/{self.sptitname}/{exp_from}/train.jsonl'))
            cat_stats = self.get_cat_stats(exp_from)
            setting_name = f'Settings : {self.bi_firt_row_text}'


            epoch_for_report = best_epoch['Epoch']+1
            # row = {setting_name:f"{exp}, ({num_samples})({cat_stats}), epoch {epoch_for_report}", 
            #     'Overall':best_epoch['Overall'], 'MRR':best_epoch['MRR']}
            row = {setting_name:f"{exp}", 
                'Overall':best_epoch['Overall'], 'MRR':best_epoch['MRR'], 'samples':num_samples}


            
            row['HO'] = best_epoch['HO']
            row['LO'] = best_epoch['LO']
            row['NO'] = best_epoch['NO']

            
            pos, neg, steps, warmup, bs = self.get_pos_stats(exp_from, bienc_model_dir, best_epoch['Epoch'])
            row['POS'] = round(pos, 2)
            row['NEG'] = round(neg, 2)
            row['STEPS'] = steps
            row['WARMUP'] = warmup
            row['BS'] = bs
            row['EPOCH'] = epoch_for_report

            row['DD'] = cat_stats
            

            for recall_i in best_epoch['recalls']:
                recall_i_data = {"Settings : ":row[setting_name]}
                for key in best_epoch['recalls'][recall_i]:
                    recall_i_data[key] = best_epoch['recalls'][recall_i][key]
                
                recall_i_data['MRR'] = row['MRR']
                recall_i_data['POS'] = row['POS']
                recall_i_data['NEG'] = row['NEG']
                recall_i_data['STEPS'] = row['STEPS']
                recall_i_data['WARMUP'] = row['WARMUP']
                recall_i_data['BS'] = row['BS']
                recall_i_data['EPOCH'] = row['EPOCH']
                recall_i_data['samples'] = row['samples']
                recall_i_data['DD'] = row['DD']

                recall_at_n_reports[recall_i].append(recall_i_data)

            
            
            
            recall_at_1_reports.append(row)


        
        df_recall1 = pd.DataFrame(recall_at_1_reports)
        title = f"{self.onto.upper()}:  recall@1"
        ppt_report.add_new_slide(df_recall1, title)
    
        for reacall_at_i in self.bi_report_start_end:
            if reacall_at_i == 'recall@1':
                continue

            df = pd.DataFrame(recall_at_n_reports[reacall_at_i])
            if reacall_at_i == 'recall@64':
                # df['MRR'] = df_recall1['MRR']
                # df['POS'] = df_recall1['POS']
                # df['NEG'] = df_recall1['NEG']
                # df['STEPS'] = df_recall1['STEPS']
                # df['WARMUP'] = df_recall1['WARMUP']
                # df['BS'] =  df_recall1['BS']
                # df['EPOCH'] =  df_recall1['EPOCH']
                df = df[['Settings : ', 'Overall', 'MRR', 'HO', 'LO', 'NO','POS', 'NEG', 'STEPS', 'WARMUP', 'BS', 'EPOCH', "DD"]]

            title = f"{self.onto.upper()}:  {reacall_at_i}"
            ppt_report.add_new_slide(df, title)


        ppt_file_path = f"{self.bi_report_save_to}/{self.onto}_blink_bienc_evaluation"
        ppt_report.save_pptx(ppt_file_path)

        recall_at_n_reports['recall@1'] = recall_at_1_reports

        with open(f'{ppt_file_path}.json', 'w') as f:
            json.dump(recall_at_n_reports, f, indent=1)

def biencoder_eval_report(params,bi_pred_file):
        with open(bi_pred_file) as f:
            bi_all_predictions = json.load(f)
        data_dir = params["data_path"]
        kb_dict = {}
        with io.open(f'{data_dir}/kb.jsonl', mode="r", encoding="utf-8") as file:
            for line in file:
                e = json.loads(line.strip())
                kb_dict[e['id']]=e
        with open(f'{data_dir}/id_map.json', 'r') as f:
            id_map = json.load(f)
        with open(f'{params["grag_data_path"]}/{params["mode"]}_grag.json') as f:
            multiple_gt_grag = json.load(f)
            multiple_gt_grag_dict = {}
            for item in multiple_gt_grag:
                multiple_gt_grag_dict[item['sample_id']] = item
        converted = []
        for bi_pred in bi_all_predictions:
            bi_pred_sample_id = bi_pred['mention_data']['sample_id']
            retriever_predictions = []
            for c in bi_pred['retriever_predictions']:
                c['title'] = kb_dict[c['id']]['title']
                c['id'] = id_map[str(c['id'])]
                retriever_predictions.append(c)

            m = bi_pred['mention_data']
            mention = m['mention']
            context_left = m['context_left'].strip()
            context_right = m['context_right'].strip()
            if context_left == '' and context_right != '':
                mention_context = '[MENTION_START] '+mention.strip()+' [MENTION_END] '+context_right
            elif context_left != '' and context_right == '':
                mention_context = context_left+' [MENTION_START] '+mention.strip()+' [MENTION_END]'
            elif context_left == '' and context_right == '':
                mention_context = '[MENTION_START] '+mention.strip()+' [MENTION_END]'
            else:
                mention_context = m['context_left'].strip() + ' [MENTION_START] '+mention.strip()+' [MENTION_END] '+m['context_right'].strip()


            d = {'sample_id' : bi_pred_sample_id,
            'mention': mention,
            'mention_context':mention_context.strip(),
            'ground_truth': multiple_gt_grag_dict[bi_pred_sample_id]['ground_truth'],
            'retriever_result_gt':bi_pred['retriever_retrived_gt'],
            'retriever_predictions':retriever_predictions
            }
            converted.append(d)

        # Unlabeled splits (the QTL train corpora ship with ground_truth == [])
        # have nothing to score against: MultiGTEvaluation would build an empty
        # not_none_data and calculate_mrr would divide by zero. Write a note
        # instead of metrics. This is per-split and data-driven on purpose --
        # config's has_gt is per-corpus, so it would also skip the labeled
        # test split of the same corpus.
        if not any(d['ground_truth'] for d in converted):
            report = (f'Bi-Encoder\n{"_"*20}\n'
                      f'No ground truth in the "{params["mode"]}" split '
                      f'({len(converted)} mentions) - evaluation skipped.\n'
                      f'Candidates were still written to {bi_pred_file}.\n{"_"*20}')
            with open(f"{bi_pred_file.replace('.json', '_eval.txt')}", 'w') as f:
                f.write(report)
            print(report)
            return

        with open(f'{params["kb_file_path"]}') as f:
            exact_kb = json.load(f)
        eval_bi = MultiGTEvaluation(converted, exact_kb, 'retriever_predictions', multiple_gt_grag, for_retrieval=True)
        not_none_data, none_data = eval_bi.get_report()
        report = f'Bi-Encoder\n{"_"*20}\n{eval_bi.text_report}\n{"_"*20}'
        with open(f"{bi_pred_file.replace('.json', '_eval.txt')}", 'w') as f:
            f.write(report)


