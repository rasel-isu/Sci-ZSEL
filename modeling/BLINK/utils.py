import io
import json
import re
import unicodedata
import networkx as nx
from typing import Dict, Set, List, Any, Optional
import pandas as pd
import matplotlib as mpl
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches
import torch

with open('../../config.json') as f:
    CONFIG = json.load(f)


UNICODE_MAP = {

    "\u2013": "-",   # en dash
    "\u2014": "-",   # em dash
    "\u2212": "-",   # minus sign
    "\u223c": "~",   # tilde operator
    "\u2019": "'",   # right single quote
    "\u2018": "'",   # left single quote
    "\u201c": '"',   # left double quote
    "\u201d": '"',   # right double quote
    "\u00b7": "*",   # middle dot (often used as multiplication in papers)

    # Math / comparison operators
    '×': 'x',          # cross (breed crosses: Meishan x Landrace)
    '±': '+/-',
    '≤': '<=',
    '≥': '>=',
    '≠': '!=',
    '≈': '~',
    '∼': '~',
    '→': '->',
    '−': '-',           # minus sign (not hyphen)
    '⁻': '-',           # superscript minus
    '﹥': '>',           # small greater-than

    # Superscripts (e.g. 10⁵, 10⁶)
    '⁵': '5',
    '⁶': '6',

    # Greek letters (used in statistical/biological notation)
    'α': 'alpha',
    'β': 'beta',
    'γ': 'gamma',
    'ɣ': 'gamma',
    'μ': 'mu',

    # Punctuation variants
    '–': '-',           # en dash
    "'": "'",           # curly apostrophe
    '·': '.',           # middle dot

    # Accented Latin (author names, breed names)
    'é': 'e',
    'ê': 'e',
    'ö': 'o',
    'ł': 'l',

    # Other
    '®': '',            # registered trademark — just remove
    '°': ' degrees',
    'І': 'I',           # Cyrillic I that looks like Latin I
}
class DataCleaner:
    def __init__(self):
        pass
    def fix_spaced_numbers(self, text):
        # Fix "0. 05" → "0.05", "8. 4" → "8.4"
        text = re.sub(r'(\d)\. (\d)', r'\1.\2', text)
        # Fix "1, 439" → "1,439"
        text = re.sub(r'(\d), (\d)', r'\1,\2', text)
        return text
    
    def normalize_unicode(self, text):
        for char, replacement in UNICODE_MAP.items():
            text = text.replace(char, replacement)
        text = unicodedata.normalize('NFKC', text)
        return text
    
    def clean_unicode(self, s: str) -> str:
        if not s:
            return s
        # NFD splits chars into base + combining marks; drop stray combining marks
        # that aren't attached to a letter (the U+0301 case from PDF extraction).
        out = []
        for ch in unicodedata.normalize("NFD", s):
            if unicodedata.category(ch) == "Mn":
                # Skip orphan combining mark only if previous char isn't a letter
                if out and unicodedata.category(out[-1]).startswith("L"):
                    out.append(ch)
                # else drop it
            else:
                out.append(ch)
        s = unicodedata.normalize("NFC", "".join(out))
        # Apply lookalike replacements
        for k, v in UNICODE_MAP.items():
            s = s.replace(k, v)
        # Collapse any double spaces introduced
        while "  " in s:
            s = s.replace("  ", " ")
        return s.strip()
    
    def normalize_whitespace(self, text):
        text = re.sub(r' +', ' ', text)   # collapse multiple spaces
        text = text.strip()
        return text

    def clean_copid_text_from_pdf(self, text):
        # Join hyphenated line breaks: "hyper-\ntrophy" -> "hypertrophy"
        text = re.sub(r"-\n(\w)", r"\1", text)
        # Replace remaining newlines with single space
        text = text.replace("\r", " ").replace("\n", " ")
        text = self.clean_unicode(text)
        return text.strip()
    
    def clean_mention_context(self, text):
        text = self.fix_spaced_numbers(text)
        text = self.normalize_unicode(text)
        text = self.normalize_whitespace(text)
        return text


class PPTReport:
    def __init__(self):
        self.presentation = Presentation()

    def save_pptx(self, file_path):
        self.presentation.save(f"{file_path}.pptx")



    
    def add_new_slide(self, df, title):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[5])
        slide.shapes.title.text = title
        x, y, cx, cy = Inches(1), Inches(1.6), Inches(8), Inches(2)
        rows, cols = df.shape[0] + 1, df.shape[1]
        table = slide.shapes.add_table(rows, cols, x, y, cx, cy).table

        for table_col in table.columns:
            table_col.width = Inches(1)

        table.columns[0].width = Inches(2)

        for col_idx, column_name in enumerate(df.columns):
            table.cell(0, col_idx).text = column_name

        for row_idx, row in df.iterrows():
            for col_idx, value in enumerate(row):
                table.cell(row_idx + 1, col_idx).text = str(value)


def get_word_overlap_count(term_1, term_2):
    words_term_1  = set(term_1.split())
    words_term_2 = set(term_2.split())
    common_words = words_term_1.intersection(words_term_2)
    return len(common_words)

def get_term_similarity(model, term_1, term_2, word_overlap=False):
    if word_overlap:
        return get_word_overlap_count(term_1, term_2)
    else:
        sim_score = round(model.similarity(model.encode(term_1), model.encode(term_2)).item(), 2)
        return sim_score
        
def does_ent_sim_smaller_than_its_parent_child_or_threshhold(model, item_id, relations, kb, kb_prime_def_m, ent_prime):
    ent_old_name = kb_prime_def_m[item_id]['old_name']
    # ent_prime = kb_prime_def_m[item_id]['newly_generated_name']
    ent_score = get_term_similarity(model, ent_old_name, ent_prime)

    score_dict = {}
    scores = []
    for relation in relations:
        if item_id in relation:
            ent_relation = relation[item_id]
            if 'parents' in ent_relation:
                ent_parents = ent_relation['parents']
                for parent_ent_id in ent_parents:
                    parent_ent = kb[parent_ent_id]
                    parent_ent_name = parent_ent['name']
                    parent_score = get_term_similarity(model, parent_ent_name, ent_prime)
                    scores.append(parent_score)
                    score_dict[parent_ent_name]=parent_score
            if 'children' in ent_relation:
                ent_children = ent_relation['children']
                for child_ent_id in ent_children:
                    child_ent = kb[child_ent_id]
                    child_ent_name = child_ent['name']
                    child_score = get_term_similarity(model, child_ent_name, ent_prime)
                    scores.append(child_score)
                    score_dict[child_ent_name]=child_score

    ent_sim_smaller_than_its_parent_child = False
    for scr in scores:
        if ent_score<scr:
            ent_sim_smaller_than_its_parent_child = True
            break

    if ent_score < 0.9:
        ent_sim_smaller_than_its_parent_child = True


    parent_child_score = str({'ent':ent_old_name, 'score':ent_score,'prime':ent_prime, 'parent_child':score_dict}).replace("\\", '')
    
    return ent_sim_smaller_than_its_parent_child, parent_child_score


def get_mesh_graph():
    # G = nx.MultiDiGraph()
    G = nx.DiGraph()
    rfile = [
        'relations_desc2025.json',
        'relations_pa2025.json',
        'relations_qual2025.json',
        'relations_supp2025.json'
        ]
    for i in rfile:
        with open('data/bc5cdr/onto/'+i) as f:
            data = json.load(f)
        for source, relations in data.items():
            for rel_type, targets in relations.items():
                for target in targets:
                    G.add_edge(source, target, relation=rel_type)
        print(f"Total nodes: {G.number_of_nodes()}")
        print(f"Total edges: {G.number_of_edges()}")

    return G

def path_with_relations(G, source, target):

    try:
        if not nx.has_path(G, source, target):
            return [], '', -1
    except Exception as e:
        return [], '', -1

    route = ''
    depth = []
    path = nx.shortest_path(G, source=source, target=target)
    for d, (u, v) in enumerate(zip(path[:-1], path[1:]), 1):
        depth.append(d)
        edge_data = G.get_edge_data(u, v)
        edge = edge_data[0]
        route += f"{u}, {edge['relation']}, {v} | "
    hops = len(path)-1
    # if hops>1:
    #     print(hops)

    return path, route.strip(), hops

def get_hop_count(data, json_file):
    df = pd.DataFrame(data)
    counts = df['hops'].value_counts().to_dict()
    samples = []
    samples_hist = []
    for hop, count in counts.items():
        samples_hist.extend([int(hop)] * count)
        # if hop == -1:
        #     continue
        samples.extend([int(hop)] * count)

    samples_array = np.array(samples)
    hops_df = pd.DataFrame(samples_array, columns=["hops"])
    percentiles = hops_df["hops"].quantile([0.70, 0.80, 0.90]).to_dict()
    counts = dict(sorted(counts.items()))
    total = sum(counts.values())
    counts_prct = {}
    xtik_hop_label = {}
    for item in counts:
        perct = round(counts[item]/total, 2)
        pc_r = int(round(perct*100, 0))
        counts_prct[f'{item} ({perct})'] = counts[item]
        if item == -1:
            xtik_hop_label[item] = f'({counts[item]}, {pc_r}%) Not'
        else:
            xtik_hop_label[item] = f'({counts[item]}, {pc_r}%) {item}'
    
    mpl.rcParams['font.size'] = 25
    mpl.rcParams['axes.titleweight'] = 'bold'
    mpl.rcParams['axes.labelweight'] = 'bold'

    xtik_hops = range(-1, 11, 1)
    plt.figure(figsize=(10, 25))
    plt.hist(samples_hist, bins=xtik_hops, edgecolor='black', align='left', rwidth=0.8)
    # plt.title("Hop Count Distribution")
    plt.xlabel("(N. of sample,%) Hop")
    plt.ylabel("Hop frequency")
    plt.grid(True, linestyle='--', alpha=0.6)
    plt.xticks(xtik_hops, [xtik_hop_label[hop] if hop in xtik_hop_label else str(hop) for hop in xtik_hops], rotation=75)
    hist_file = json_file.replace('.json', '_hist.png')
    plt.savefig(hist_file, dpi=300)
    plt.close()  

    

def compare_all_settings_hist(onto, data, save_to_dir):
    # Prepare layout
    half_len_data = int(len(data)/2)
    fig, axes = plt.subplots(half_len_data, half_len_data, figsize=(15, 10))
    axes = axes.flatten()

    # Set common x-ticks and y-limit
    all_keys = sorted(set(k for d in data.values() for k in d))
    x_ticks = list(all_keys)
    y_max = max(max(d.values()) for d in data.values()) + 50

    # # Plot each histogram
    # for i, (title, count_dict) in enumerate(data.items()):
        
    #     ax = axes[i]
    #     total = sum(count_dict.values())
    #     counts = [count_dict.get(k, 0) for k in x_ticks]
    #     bars = ax.bar(x_ticks, counts, color="#0A5693", edgecolor='black')

    #     # Annotate bars with count and percentage
    #     for bar, c in zip(bars, counts):
    #         pct = (c / total) * 100
    #         pct = int(round(pct, 0))

    #         ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 2,
    #                 f'{c}\n{pct}%', ha='center', va='bottom', fontsize=9)

    #     ax.set_title(title, fontsize=14)
    #     ax.set_ylim(0, y_max)
    #     ax.set_xticks(x_ticks)
    #     ax.set_xlabel('Hops')
    #     ax.set_ylabel('Count')

    # plt.tight_layout()
    # plt.savefig(f"{save_to_dir}{onto}_subplots.png", dpi=500)

    methods = list(data.keys())
    all_labels = sorted(set(k for d in data.values() for k in d))  # all possible class labels
    label_strs = [str(k) for k in all_labels]

    x = np.arange(len(all_labels))  # positions for groups
    width = 0.1  # width of each bar
    # Define color palette
    colors = ['blue','orange','green','red','purple','brown','pink','gray']
    # Prepare figure
    fig, ax = plt.subplots(figsize=(18, 10))
    # Plot each method as separate bar set
    for i, (method, color) in enumerate(zip(methods, colors)):
        counts = [data[method].get(k, 0) for k in all_labels]
        total = sum(counts)
        positions = x + (i - 1.5) * width
        bars = ax.bar(positions, counts, width, label=method, color=color)

        # Add count and % on top of each bar
        for bar, c in zip(bars, counts):
            pct = (c / total) * 100 if total > 0 else 0
            pct = int(round(pct, 0))
            ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 2,
                    f'      {c}, {pct}%', ha='center', va='bottom', fontsize=9, rotation=90)

    # Formatting
    ax.set_xticks(x)
    ax.set_xticklabels(label_strs)
    ax.set_xlabel('Hops')
    ax.set_ylabel('No. of sample')
    ax.set_title('Hop distribution')
    ax.legend()
    ax.set_ylim(0, max(max(data[m].get(k, 0) for k in all_labels) for m in methods) + 50)
    ax.grid(True, axis='y', linestyle='--', linewidth=0.5, alpha=0.7)
    plt.tight_layout()
    plt.savefig(f"{save_to_dir}{onto}_single_hist.png", dpi=500)

def get_mesh_relations():
    rfile = ['relations_desc2025.json',
        'relations_pa2025.json',
        'relations_qual2025.json',
        'relations_supp2025.json']
    relations=[]
    for i in rfile:
        with open('data/bc5cdr/onto/'+i) as f:
            data = json.load(f)
        relations.append(data)
    return relations

def get_medic_relations():
    rfile = ['only_medic_def.json']
    relations=[]
    for i in rfile:
        with open('data/ncbi/onto/'+i) as f:
            data = json.load(f)
            relation_dict = {}
            for ent in data:
                if isinstance( data[ent]['ParentIDs'], list):
                    ent_rel = [ en.replace('MESH:', '') for en in data[ent]['ParentIDs'] ]
                else:
                    if data[ent]['ParentIDs'] == "MESH:C":
                        continue

                    ent_rel = [ data[ent]['ParentIDs'].replace('MESH:', '') ]

                relation_dict[ent] = {'parents' : ent_rel}
            
            relations.append(relation_dict)

    return relations

def get_animal_science_relations(onto):
    rfile = [f'{onto}_kb.json']
    relations=[]
    for i in rfile:
        with open(f'data/{onto}/onto/'+i) as f:
            data = json.load(f)
            relation_dict = {}
            for ent in data:
                if isinstance( data[ent]['ParentIDs'], list):
                    ent_rel = [ en.replace('MESH:', '') for en in data[ent]['ParentIDs'] ]
                else:
                    if data[ent]['ParentIDs'] == "MESH:C":
                        continue

                    ent_rel = [ data[ent]['ParentIDs'].replace('MESH:', '') ]

                relation_dict[ent] = {'parents' : ent_rel}
            
            relations.append(relation_dict)

    return relations


def get_relations_from_ontology(onto, kb_filepath):
    rfile = [kb_filepath]
    relations=[]
    for i in rfile:
        with open(i) as f:
            data = json.load(f)
            relation_dict = {}
            for ent in data:
                if isinstance( data[ent]['ParentIDs'], list):
                    ent_rel = [ en.replace('MESH:', '') for en in data[ent]['ParentIDs'] ]
                else:
                    if data[ent]['ParentIDs'] == "MESH:C":
                        continue

                    ent_rel = [ data[ent]['ParentIDs'].replace('MESH:', '') ]

                relation_dict[ent] = {'parents' : ent_rel}
            
            relations.append(relation_dict)

    return relations

def category_specific_count(mention, title):
    mention_lower = mention.lower()
    title_lower = title.lower()
    
   
    if mention_lower == title_lower:
        return "HO"
    elif mention_lower in title_lower and title_lower != mention_lower:
        return "MINT"
    else:
        lo = False
        mention_words = mention_lower.split()
        for word in mention_words:
            if word in title_lower:
                lo = True
                break
        if lo:
            return "LO"
        else:
            return "NO"
        
def get_category(mention, title):

    mention_lower = mention.lower().replace(',', ' ')
    if title is None:
        title = ''
    title_lower = title.lower().replace(',', ' ')
    subcat = check_plural(mention_lower, title_lower)

    if mention_lower == title_lower:
        return 'HO'
    elif subcat=='plural':
        return 'HO'
    elif mention_lower in title_lower and title_lower != mention_lower:
        return 'MINT'
    else:
        words_mention = set(mention_lower.split())
        words_title = set(title_lower.split())
        common_words = words_mention.intersection(words_title)
        if common_words:
            return "LO"
        else:
            for word in words_mention:
                if f' {word} '  in f' {title_lower} ':
                    return "LO"

            return "NO"
        



def get_category_test(mention, title):
    mention_lower = mention.lower()
    title_lower = title.lower()
    if mention_lower == title_lower:
        return 'HO'
    elif mention_lower in title_lower and title_lower != mention_lower:
        return 'MINT'
    else:
        words_mention = set(mention_lower.split())
        words_title = set(title_lower.split())
        common_words = words_mention.intersection(words_title)
        if common_words:
            return "LO"
        else:
            return "NO"
        
def check_plural(mention, title):
    if title is None:
        title = ''
    mention, title = mention.lower(), title.lower()
    # if len(mention) != len(gt_title):
    longer, shorter = (mention, title) if len(mention) > len(title) else (title, mention)
    diff = len(longer)-len(shorter)
    if diff ==1 and longer.startswith(shorter) and longer.endswith("s"):
        return 'plural'
    else:
        return 'Pure'
    
def read_jsonl(filename):
    data = []
    if '.jsonl' in filename:
        with io.open(filename, mode="r", encoding="utf-8") as file:
            for line in file:
                data.append(json.loads(line.strip()))
    return data

def read_json(p):
    with open(p) as f:
        return json.load(f)
    
def compare_with_multiple_gt(candidate, gt_list):
    for gt_item in gt_list:
        if gt_item['id'] == candidate['id']:
            return True, gt_item
    return False, {}
    
        
def compare_with_multiple_gt_with_altid(kb, candidate, gt_list):
    candidate_id = candidate['id']
    if candidate_id != '-1':
        if 'altdiseaseid' in kb[candidate_id]:
            altids = kb[candidate_id]['altdiseaseid']
            for alid in altids: 
                for gt_item in gt_list:
                    if gt_item['id'] == alid:
                        return True, gt_item,  alid
                
    return False, {}, None

def is_accuratly_pseudo_labelled(gt_id_list, pseudo_ent):

    
    if 'altdiseaseid' in pseudo_ent:
        pseudo_label_id_list = [pseudo_ent["id"]] + pseudo_ent["altdiseaseid"]
    else:
        pseudo_label_id_list = [pseudo_ent["id"]]

    for pseudo_label_id in pseudo_label_id_list:
        if pseudo_label_id in gt_id_list:
            return True
                
def gt_mention_category_count(source_dir, filename):
    data = read_jsonl(source_dir+filename)
    cat_wise_gt_matched = {
            "HO":{'count':0, 'items':[]}, 
            # "plural":{'count':0, 'items':[]}, 
            "MINT":{'count':0, 'items':{
                'plural':{'count':0, 'items':[]},
                'Pure':{'count':0, 'items':[]}} 
            },
            "LO":{'count':0, 'items':{
                'plural':{'count':0, 'items':[]},
                'Pure':{'count':0, 'items':[]}} 
            },
            "NO":{'count':0, 'items':{
                'plural':{'count':0, 'items':[]},
                'Pure':{'count':0, 'items':[]}} 
            }
            }
    unq_mention = {}
    unq_gt = {}
    for i in data:
        try:
            mention = i['mention']
        except Exception as e:
            print(i)
            input('s')
        title = i['label_title']
        if mention in unq_mention:
            unq_mention[mention]+=1
        else:
            unq_mention[mention]=1

        if title in unq_gt:
            unq_gt[title]+=1
        else:
            unq_gt[title]=1

        # if title == None:
        #     title = ''
            
        # if mention == 'autosomal dominant disease':
        #     print(0)
        category = get_category(mention, title)
        # category_old = category_specific_count(mention, title)
        # if category!=category_old:
        #     print(0)
        
        item_d = {
                'mention' : mention,
                'gt_title' : title,
                'gt_id' : i['label_id']
            }
        if category=="HO":
            cat_wise_gt_matched[category]['count']+=1
            cat_wise_gt_matched[category]['items'].append(item_d)
        else:
            cat_wise_gt_matched[category]['count']+=1
            subcat = check_plural(mention, title)
            cat_wise_gt_matched[category]['items'][subcat]['count']+=1
            cat_wise_gt_matched[category]['items'][subcat]['items'].append(item_d)
            
    cat_wise_gt_matched['unique_mention_count'] = len(unq_mention)
    cat_wise_gt_matched['unique_gt_count'] = len(unq_gt)

    sorted_unq_mention = sorted(unq_mention.items(), key=lambda item: item[1], reverse=True)
    cat_wise_gt_matched['unique_mention'] =  dict(sorted_unq_mention)

    sorted_unq_gt = sorted(unq_gt.items(), key=lambda item: item[1], reverse=True)
    cat_wise_gt_matched['unique_gt'] =  dict(sorted_unq_gt)

    with open(source_dir+filename.replace('.jsonl', '_category_count.json'), 'w') as f:
        json.dump(cat_wise_gt_matched, f, indent=1)

    return cat_wise_gt_matched



class MEDICGraph:
    def __init__(self, json_path: str):
        self.graph = nx.DiGraph()  # Directed graph: child -> parent
        self.entities = {}
        self._load_data(json_path)
        print(f"Graph created: {self.graph.number_of_nodes()} nodes, {self.graph.number_of_edges()} edges")
    
    def _normalize_parent_ids(self, parent_ids: Any) -> List[str]:
        if isinstance(parent_ids, str):
            parent = parent_ids.replace('MESH:', '')
            return [parent]
        elif isinstance(parent_ids, list):
            parents = [ i.replace('MESH:', '') for i in parent_ids]
            return parents
        return []
    
    def _load_data(self, json_path: str):
        with open(json_path, 'r', encoding='utf-8') as f:
            self.entities = json.load(f)
        print(f"Loaded {len(self.entities)} entities")
        for entity_id in self.entities:
            entity_data = self.entities[entity_id]
            self.graph.add_node(entity_id, **entity_data)
        for entity_id in self.entities:
            entity_data = self.entities[entity_id]
            parent_ids = self._normalize_parent_ids(entity_data.get('ParentIDs', []))
            for parent_id in parent_ids:
                # Add edge: child -> parent
                self.graph.add_edge(entity_id, parent_id, relationship='ParentIDs')
    
    def get_connected_ids(self, entity_id: str, depth: int) -> Dict[int, Set[str]]:
        if entity_id not in self.graph:
            raise ValueError(f"Entity ID '{entity_id}' not found in graph")
        
        result = {0: {entity_id}}
        
        if depth == 0:
            return result
        
        undirected = self.graph.to_undirected()
        
        visited = {entity_id}
        current_level = {entity_id}
        
        for hop in range(1, depth + 1):
            next_level = set()
            
            for node in current_level:
                # Get all neighbors (both directions)
                neighbors = set(undirected.neighbors(node))
                
                # Add unvisited neighbors
                for neighbor in neighbors:
                    if neighbor not in visited:
                        next_level.add(neighbor)
                        visited.add(neighbor)
            
            result[hop] = next_level
            current_level = next_level
            
            if not current_level:
                break
        
        return result
    
    def get_family_relationships(self, entity_id: str) -> Dict[str, Set[str]]:
        """
        Get immediate family relationships: parents, children, and siblings.
        
        Args:
            entity_id: The entity ID to query
            
        Returns:
            Dictionary with:
            - 'parents': Set of parent IDs (successors in directed graph)
            - 'children': Set of child IDs (predecessors in directed graph)
            - 'siblings': Set of sibling IDs (share at least one parent)
        """
        if entity_id not in self.graph:
            print(f"Entity ID '{entity_id}' not found in graph")
            return {
                'parents': set(),
                'children': set(),
                'siblings': set()
            }
        
        # Parents are successors (we go from child -> parent)
        parents = set(self.graph.successors(entity_id))
        
        # Children are predecessors (who points to this node)
        children = set(self.graph.predecessors(entity_id))
        
        # Siblings: all children of this entity's parents, excluding itself
        siblings = set()
        for parent_id in parents:
            # Get all children of this parent (predecessors)
            siblings.update(self.graph.predecessors(parent_id))
        siblings.discard(entity_id)  # Remove self
        
        return {
            'parents': parents,
            'children': children,
            'siblings': siblings
        }
    
    def get_entity_info(self, entity_id: str) -> Optional[Dict[str, Any]]:
        """Get full information about an entity."""
        return self.entities.get(entity_id)
    
    def get_ancestors(self, entity_id: str, max_depth: int = 100) -> Dict[int, Set[str]]:
        """
        Get all ancestors (parents, grandparents, etc.).
        
        Args:
            entity_id: Starting entity ID
            max_depth: Maximum depth to traverse upward
            
        Returns:
            Dictionary mapping hop level to set of ancestor IDs
        """
        if entity_id not in self.graph:
            raise ValueError(f"Entity ID '{entity_id}' not found in graph")
        
        result = {0: {entity_id}}
        visited = {entity_id}
        current_level = {entity_id}
        
        for hop in range(1, max_depth + 1):
            next_level = set()
            
            for node in current_level:
                # Get parents (successors)
                parents = set(self.graph.successors(node))
                
                for parent in parents:
                    if parent not in visited:
                        next_level.add(parent)
                        visited.add(parent)
            
            result[hop] = next_level
            current_level = next_level
            
            if not current_level:
                break
        
        return result
    
    def get_descendants(self, entity_id: str, max_depth: int = 100) -> Dict[int, Set[str]]:
        """
        Get all descendants (children, grandchildren, etc.).
        
        Args:
            entity_id: Starting entity ID
            max_depth: Maximum depth to traverse downward
            
        Returns:
            Dictionary mapping hop level to set of descendant IDs
        """
        if entity_id not in self.graph:
            raise ValueError(f"Entity ID '{entity_id}' not found in graph")
        
        result = {0: {entity_id}}
        visited = {entity_id}
        current_level = {entity_id}
        
        for hop in range(1, max_depth + 1):
            next_level = set()
            
            for node in current_level:
                # Get children (predecessors)
                children = set(self.graph.predecessors(node))
                
                for child in children:
                    if child not in visited:
                        next_level.add(child)
                        visited.add(child)
            
            result[hop] = next_level
            current_level = next_level
            
            if not current_level:
                break
        
        return result
    
    def shortest_path(self, start_id: str, end_id: str) -> Optional[List[str]]:
        """
        Find shortest path between two entities.
        
        Args:
            start_id: Starting entity ID
            end_id: Target entity ID
            
        Returns:
            List of entity IDs representing the path, or None if no path exists
        """
        try:
            # Use undirected version for path finding
            undirected = self.graph.to_undirected()
            return nx.shortest_path(undirected, start_id, end_id)
        except (nx.NetworkXNoPath, nx.NodeNotFound):
            return None
    
    def get_statistics(self) -> Dict[str, Any]:
        """Get graph statistics using NetworkX functions."""
        return {
            'total_nodes': self.graph.number_of_nodes(),
            'total_edges': self.graph.number_of_edges(),
            'avg_degree': sum(dict(self.graph.degree()).values()) / self.graph.number_of_nodes(),
            'avg_in_degree': sum(dict(self.graph.in_degree()).values()) / self.graph.number_of_nodes(),
            'avg_out_degree': sum(dict(self.graph.out_degree()).values()) / self.graph.number_of_nodes(),
            'is_directed': self.graph.is_directed(),
            'is_connected': nx.is_weakly_connected(self.graph)
        }
    
    def print_tree(self, entity_id: str, max_depth: int = 3, 
                   direction: str = 'down', indent: str = "", is_last: bool = True):
        """
        Print a tree view of the entity and its descendants or ancestors.
        
        Args:
            entity_id: Root entity ID
            max_depth: Maximum depth to display
            direction: 'down' for children, 'up' for parents
            indent: Current indentation
            is_last: Whether this is the last child (for formatting)
        """
        if entity_id not in self.graph:
            print(f"{indent}[{entity_id} not found]")
            return
        
        entity = self.entities.get(entity_id, {})
        name = entity.get('name', 'N/A')
        
        # Print current node
        if indent:
            connector = "└── " if is_last else "├── "
            print(f"{indent}{connector}{entity_id}: {name}")
        else:
            print(f"{entity_id}: {name}")
        
        if max_depth <= 0:
            return
        
        # Get connected nodes
        if direction == 'down':
            connected = list(self.graph.predecessors(entity_id))
        else:
            connected = list(self.graph.successors(entity_id))
        
        # Print children/parents
        for i, connected_id in enumerate(sorted(connected)):
            is_last_child = (i == len(connected) - 1)
            new_indent = indent + ("    " if is_last else "│   ")
            self.print_tree(connected_id, max_depth - 1, direction, new_indent, is_last_child)

class MESHGraph(MEDICGraph):
    def __init__(self):
        self.graph = self.get_mesh_graph()
        print(f"Graph created: {self.graph.number_of_nodes()} nodes, {self.graph.number_of_edges()} edges")

    def get_mesh_graph(self):
        G = nx.DiGraph()
        rfile = [
            'relations_desc2025.json',
            'relations_pa2025.json',
            'relations_qual2025.json',
            'relations_supp2025.json'
            ]
        for i in rfile:
            with open('data/bc5cdr/onto/'+i) as f:
                data = json.load(f)
            for source, relations in data.items():
                for rel_type, targets in relations.items():
                    if rel_type in ['parents', 'children']:
                        for target in targets:
                            G.add_edge(source, target, relation=rel_type)
        return G


    
    

def get_connected_ents_for_the_label(graph_obj, map_int_to_kb, 
                       map_kb_to_int, ent_int_id, include_siblings=False):
    ent_kb_id = map_int_to_kb[str(ent_int_id)]
    family = graph_obj.get_family_relationships(ent_kb_id)
    selected_ents = list(family['parents']) + list(family['children'])
    if include_siblings:
        selected_ents+=list(family['siblings'])
    selected_ents = list(set(selected_ents))
    # if len(selected_ents) <= 20:
    #     print(selected_ents)

    selected_ent_int_id = []
    for ent in selected_ents:
        ent_cleaned = ent.replace('MESH:', '')
        if ent_cleaned in map_kb_to_int:
            sel_ent_id = int(map_kb_to_int[ent_cleaned]) 
            selected_ent_int_id.append(sel_ent_id)

    return selected_ent_int_id


def see_candidates_and_labels(params, tokenizer, candidates_batch, 
                              original_label, sample_id):
    
    
    with open(f'{params["raw_data_path"]}/id_map.json') as f:
        id_map  = json.load(f)
    train_data_dict = {}
    for sample in read_jsonl(f'{params["raw_data_path"]}/train.jsonl'):
        train_data_dict[sample['sample_id']] = sample
    
    the_sample = train_data_dict[int(sample_id)]
    label_id = id_map[str(the_sample['label_id'])]
    if params['onto'] == 'ncbi':
        with open('data/ncbi/onto/only_medic_def.json') as f:
            onto  = json.load(f)
    ParentIDs = onto[label_id]['ParentIDs']
    
    decoded_candidates = f'SAMPLE ID : {sample_id}\nLABEL : {label_id}\nParentIDs : {ParentIDs}\n{"*"*100}\n\n'
    for i, c in enumerate(candidates_batch):
        decoded = tokenizer.decode(c, skip_special_tokens=False)
        if i == original_label:
            decoded_candidates+=f'{"*"*10}\n\n LABEL {"*"*10}\n\n'
        decoded_candidates+=f'{decoded}\n'
        
    with open('test.txt', 'w') as f:
        f.write(decoded_candidates)



def debug_find_problematic_sample(text_vecs, context_len, score_fn):
    batch_size = text_vecs.shape[0]
    num_candidates = text_vecs.shape[1]
    
    problematic_samples = []
    
    print(f"Testing {batch_size} samples individually...")
    
    for i in range(batch_size):
        try:
            # Test single sample
            single_sample = text_vecs[i:i+1]  # [1, num_candidates, embd_dim]
            scores = score_fn(single_sample, context_len)
            print(f"  Sample {i}: OK")
        except Exception as e:
            print(f"  Sample {i}: ERROR - {str(e)[:100]}")
            problematic_samples.append(i)
            
            # Print details about this sample
            print('error item \n\n')
            # print(text_vecs)
            # print(f"    Shape: {text_vecs[i].shape}")
            # print(f"    Min: {text_vecs[i].min()}, Max: {text_vecs[i].max()}")
            # print(f"    Has NaN: {torch.isnan(text_vecs[i]).any()}")
            # print(f"    Has Inf: {torch.isinf(text_vecs[i]).any()}")
    
    return problematic_samples

def plot_train_test_accuracy(acc_and_epoch, save_to):
    train_acc = [epoch['train']['acc'] for epoch in acc_and_epoch]
    test_acc = [epoch['test']['acc'] for epoch in acc_and_epoch]
    epochs = list(range(1, len(acc_and_epoch) + 1))
    plt.figure(figsize=(10, 6))
    plt.plot(epochs, train_acc, 'b-o', label='Train Accuracy', linewidth=2, markersize=6)
    plt.plot(epochs, test_acc, 'r-s', label='Test Accuracy', linewidth=2, markersize=6)
    plt.xlabel('Epoch', fontsize=12)
    plt.ylabel('Accuracy', fontsize=12)
    plt.title('Train vs Test Accuracy', fontsize=14, fontweight='bold')
    plt.legend(fontsize=11)
    plt.grid(True, alpha=0.3)
    plt.tight_layout()
    plt.savefig(save_to, dpi=300, bbox_inches='tight')
