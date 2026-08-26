import json
import random
from copy import deepcopy
import re
from pptx import Presentation
from pptx.util import Inches
import pandas as pd
from data_preparation.utils import compare_with_multiple_gt, compare_with_multiple_gt_with_altid, get_category
import wandb

class Evaluation:
    def __init__(self, eval_data, error_analysis=False, for_retrieval=True):
        data = deepcopy(eval_data)
        self.data = data
        self.error_analysis = error_analysis
        self.for_retrieval = for_retrieval
        self.final_data = []
        self.mention_matched_at_k = {}
        self.cat_wise_gt_matched = {}
        self.mention_wise_triple_matched = {}   
        self.text_report = "\n\n\n"
        for i in data:
            if 'retrieved_candidates' not in i:
                continue
            self.final_data.append(i)
        
        removed_count = len(data) - len(self.final_data) 
        if removed_count > 0:
            self.text_report+=f"Removed {removed_count} items since they did not have retrieved_candidates key\n"
        
        self.none_data = [] 
        self.not_none_data = []
        for item in self.data:
            ground_truth_id = item["ground_truth"]["id"]
            if ground_truth_id.lower() == "none":
                self.none_data.append(item) 
            else:
                self.not_none_data.append(item)
        self.text_report+=f"Number of items with ground truth 'None': {len(self.none_data)}\n"
        self.text_report+=f"Number of items with ground truth not 'None': {len(self.not_none_data)}\n"

    def get_report(self):
        mrr = self.calculate_mrr()
        recall_at_1 = self.calculate_recall_at_k_for_graph_retriever(1)
        if not self.for_retrieval:
            recall_at_5 = self.calculate_recall_at_k_for_graph_retriever(5)
            recall_at_10 = self.calculate_recall_at_k_for_graph_retriever(10)
        # none_accuracy = self.get_none_accuracy()
        if self.not_none_data:
            retrieved_candidates = self.not_none_data[0]["retrieved_candidates"]
            if len(retrieved_candidates)>=63:
                recall_at_63 = self.calculate_recall_at_k_for_graph_retriever(63)


        return self.not_none_data, self.none_data

                
    
    def get_none_accuracy(self):
        data = self.none_data
        correct_count = 0
        for item in data:
            retrieved_candidates = item["retrieved_candidates"]
            if len(retrieved_candidates) == 0:
                correct_count += 1
            else:
                try:
                    for i in retrieved_candidates:
                        if i["id"].lower() == "none" or not i["id"]:
                            correct_count += 1
                            break
                except Exception as e:
                    if isinstance(retrieved_candidates[0], str):
                        if retrieved_candidates[0].lower() == "none":
                            correct_count += 1
                        
        accuracy = correct_count / len(data)
        self.text_report+=f"Total {len(data)} items was actually NONE, among them {correct_count} items correctly predicted as NONE.\nSo, accuracy for items with ground truth 'None' : {accuracy}\n"
        return accuracy
    

    def calculate_mrr(self):
        data = self.not_none_data
        reciprocal_ranks = []
        for item in data:
            ground_truth_id = item["ground_truth"]["id"]
            retrieved_candidates = item["retrieved_candidates"]
            
            # Find the rank of the first relevant item
            rank = 0
            for idx, candidate in enumerate(retrieved_candidates, start=1):
                if candidate["id"] == ground_truth_id:
                    rank = idx
                    break
            
            # Calculate reciprocal rank; if not found, reciprocal rank is 0
            reciprocal_rank = 1 / rank if rank > 0 else 0
            reciprocal_ranks.append(reciprocal_rank)
        
        # Calculate MRR by averaging all reciprocal ranks
        try:
            mrr = sum(reciprocal_ranks) / len(reciprocal_ranks)
            self.text_report+=f"\n\nCalculated MRR for {len(data)} items. \nSo, MRR : {mrr}\n\n"
        except ZeroDivisionError:
            mrr = 0
            self.text_report+=f"\n\nCalculated MRR for {len(data)} items. \nSo, MRR : 0 (ZeroDivisionError)\n\n"

        return mrr

    def calculate_recall_at_k(self, k):
        data = self.not_none_data
        relevant_count = 0
        for item in data:
            ground_truth_id = item["ground_truth"]["id"]
            retrieved_candidates = item["retrieved_candidates"][:k]  # Consider top-k candidates
            
            # Check if the relevant item is within the top-k candidates
            if any(candidate["id"] == ground_truth_id for candidate in retrieved_candidates):
                relevant_count += 1
        
        # Calculate recall@k
        try:
            recall_at_k = relevant_count / len(data)
            self.text_report+=f"\n\nFound {relevant_count} items out of {len(data)} within top-{k} candidates \nSo, recall@{k} : {relevant_count}/{len(data)}={recall_at_k}\n\n"
        except ZeroDivisionError:
            self.text_report+=f"\n\nFound {relevant_count} items out of {len(data)} within top-{k} candidates \nSo, recall@{k} : 0 (ZeroDivisionError) \n\n"
            recall_at_k = 0

        return recall_at_k
    

    def category_specific_count(self, item, is_matched, k):
        mention_id=item["mention_id"]
        mention=item["mention"]
        title= item["ground_truth"]['title']
        candidates = ' | '.join([i['title'] for i in item["retrieved_candidates"]])
        
        mention_lower = mention.lower()
        title_lower = title.lower()
        
        if mention_lower == title_lower:
            self.cat_wise_gt_matched[k]["HO"]['count'] += 1
            self.cat_wise_gt_matched[k]["HO"]['items'].append({
                'mention_id':mention_id, 
                'mention':mention, 
                'gt_title':title, 
                'matched': is_matched,
                'retrieved_candidates': candidates
                })
            if is_matched:
                self.cat_wise_gt_matched[k]["HO"]['matched'] += 1
        elif mention_lower in title_lower and title_lower != mention_lower:
            self.cat_wise_gt_matched[k]["MINT"]['count'] += 1
            self.cat_wise_gt_matched[k]["MINT"]['items'].append({
                'mention_id':mention_id, 
                'mention':mention, 
                'gt_title':title, 
                'matched': is_matched,
                'retrieved_candidates': candidates
                })
            if is_matched:
                self.cat_wise_gt_matched[k]["MINT"]['matched'] += 1
        else:
            words_mention = set(mention_lower.split())
            words_title = set(title_lower.split())
            common_words = words_mention.intersection(words_title)
            if common_words:
                self.cat_wise_gt_matched[k]["LO"]['count'] += 1
                self.cat_wise_gt_matched[k]["LO"]['items'].append({
                'mention_id':mention_id, 
                'mention':mention, 
                'gt_title':title, 
                'matched': is_matched,
                'retrieved_candidates': candidates
                })
                if is_matched:
                    self.cat_wise_gt_matched[k]["LO"]['matched'] += 1
            else:
                self.cat_wise_gt_matched[k]["NO"]['count'] += 1
                self.cat_wise_gt_matched[k]["NO"]['items'].append({
                'mention_id':mention_id, 
                'mention':mention, 
                'gt_title':title, 
                'matched': is_matched,
                'retrieved_candidates': candidates
                })
                if is_matched:
                    self.cat_wise_gt_matched[k]["NO"]['matched'] += 1


    
    def calculate_recall_at_k_for_graph_retriever(self, k):
        data = self.not_none_data
        self.gt_matched = []
        self.gt_did_not_matched = []

        self.cat_wise_gt_matched[k] = {
            "HO":{'count':0, 'matched':0, 'items':[]}, 
            "MINT":{'count':0, 'matched':0, 'items':[]}, 
            "LO":{'count':0, 'matched':0, 'items':[]},
            "NO":{'count':0, 'matched':0, 'items':[]}
        }
        
        self.mention_matched_at_k[k] = []
        
        relevant_count = 0

        if k==10:
            self.mention_wise_triple_matched[k]=[]

        for item in data:
            ground_truth_id = item["ground_truth"]["id"]
            
            if not self.for_retrieval:
                unique_triple = {}
                for id in item["retrieved_candidates"]:
                    unique_triple[id['id']] = id
                item["unique_triple"] = unique_triple

            retrieved_candidates = list(item["unique_triple"].keys())
            top_k = retrieved_candidates[:k] # Consider top-k candidates
            
            is_matched = False
            # Check if the relevant item is within the top-k candidates
            for i, candidate in enumerate(top_k):
                if ground_truth_id in candidate:
                    relevant_count += 1
                    is_matched = True
                    self.mention_matched_at_k[k].append(item["mention_id"])
                    
                    matched_dict = item["unique_triple"][candidate]
                    self.gt_matched.append({
                        'mention_id':item["mention_id"],
                        'mention':item["mention"],
                        'mention_context':item["mention_context"],
                        'ground_truth':item["ground_truth"],
                        'rank':i+1,
                        'matched_triple': {**{'triple':candidate}, **matched_dict},
                        'unique_triple':item["unique_triple"]
                        })
                    break

            if k==10:
                if is_matched:
                    triple_and_aug = []
                    for i in top_k:
                        if isinstance(item["unique_triple"], dict):
                            break

                        item["unique_triple"][i]['score']
                        triple_and_aug.append({'triple':i, 'aug':item["unique_triple"][i]})

                    self.mention_wise_triple_matched[k].append(
                        {
                            'mention_id':item["mention_id"],
                            'mention':item["mention"],
                            'ground_truth':item["ground_truth"],
                            f'top_{k}':triple_and_aug,
                        }
                    )
            
            # if not self.for_retrieval:
            self.category_specific_count(item,
                        is_matched=is_matched,
                        k=k
                    )
            
            if not is_matched:
                self.gt_did_not_matched.append({
                            'mention_id':item["mention_id"],
                            'mention':item["mention"],
                            'mention_context':item["mention_context"],
                            'ground_truth':item["ground_truth"],
                            'rank':-1,
                            'matched_triple':None,
                            'unique_triple':item["unique_triple"]
                            })
        
        # Calculate recall@k
        try:
            recall_at_k = relevant_count / len(data)
            self.text_report+=f"\n\nFound {relevant_count} items out of {len(data)} within top-{k} candidates \nSo, recall@{k} : {relevant_count}/{len(data)}={recall_at_k}\n\n"
        except ZeroDivisionError:
            recall_at_k = 0
            self.text_report+=f"\n\nFound {relevant_count} items out of {len(data)} within top-{k} candidates \nSo, recall@{k} : 0 (ZeroDivisionError)\n\n"
        
        # self.text_report+=f"\n\n"
        for cat in self.cat_wise_gt_matched[k]:
            count = self.cat_wise_gt_matched[k][cat]['count']
            matched = self.cat_wise_gt_matched[k][cat]['matched']
            score = matched / count if count > 0 else 0
            self.text_report+=f"{cat}, count: {count}, matched: {matched}, score: {score}\n"

        return recall_at_k

class MultiGTEvaluation:
    def __init__(self, eval_data,kb, candidate_key, multiple_gt_grag=None, error_analysis=False, for_retrieval=True):
        data = deepcopy(eval_data)
        self.data = data
        self.kb = kb
        self.kb_id_list = []
        for entid in self.kb:
            self.kb_id_list.append(entid)
            ent = self.kb[entid]
            if 'altdiseaseid' in ent:
                self.kb_id_list+=ent['altdiseaseid']
            
        self.candidate_key = candidate_key
        self.error_analysis = error_analysis
        self.for_retrieval = for_retrieval
        self.final_data = []
        self.mention_matched_at_k = {}
        self.cat_wise_gt_matched = {}
        self.mention_wise_triple_matched = {}   
        self.text_report = "\n\n\n"

        self.multiple_gt_grag_dict = {}
        if multiple_gt_grag:
            for item in multiple_gt_grag:
                 self.multiple_gt_grag_dict[item['sample_id']] = item

        for i in data:
            if 'retrieved_candidates' not in i:
                continue
            self.final_data.append(i)
        
        removed_count = len(data) - len(self.final_data) 
        if removed_count > 0:
            self.text_report+=f"Removed {removed_count} items since they did not have retrieved_candidates key\n"
        
        self.none_data = [] 
        self.not_none_data = []
        for item in self.data:
            ground_truth_id = item["ground_truth"]
            if ground_truth_id:
                self.not_none_data.append(item)
            else:
                self.none_data.append(item)

        self.text_report+=f"Number of items with ground truth 'None': {len(self.none_data)}\n"
        self.text_report+=f"Number of items with ground truth not 'None': {len(self.not_none_data)}\n"

    def get_report(self):
        self.mrr = self.calculate_mrr()
        self.recall_at_1 = self.calculate_recall_at_k_for_graph_retriever(1)
        self.recall_at_5 = self.calculate_recall_at_k_for_graph_retriever(5)
        # self.recall_at_10 = self.calculate_recall_at_k_for_graph_retriever(10)
        self.recall_at_10 = self.calculate_recall_at_k_for_graph_retriever(64)
        # none_accuracy = self.get_none_accuracy()

        return self.not_none_data, self.none_data

                
    
    def get_none_accuracy(self):
        data = self.none_data
        correct_count = 0
        for item in data:
            retrieved_candidates = item["retrieved_candidates"]
            if len(retrieved_candidates) == 0:
                correct_count += 1
            else:
                try:
                    for i in retrieved_candidates:
                        if i["id"].lower() == "none" or not i["id"]:
                            correct_count += 1
                            break
                except Exception as e:
                    if isinstance(retrieved_candidates[0], str):
                        if retrieved_candidates[0].lower() == "none":
                            correct_count += 1
                        
        accuracy = correct_count / len(data)
        self.text_report+=f"Total {len(data)} items was actually NONE, among them {correct_count} items correctly predicted as NONE.\nSo, accuracy for items with ground truth 'None' : {accuracy}\n"
        return accuracy
    

    def calculate_mrr(self):
        data = self.not_none_data
        reciprocal_ranks = []
        for item in data:
            if self.multiple_gt_grag_dict:
                gt_list = self.multiple_gt_grag_dict[item['sample_id']]['ground_truth']
                rank = 0
                for idx, candidate in enumerate(item[self.candidate_key], start=1):
                    got_ranked = False
                    got_ranked, gt_item = compare_with_multiple_gt(candidate, gt_list)
                    if got_ranked:
                        rank = idx
                        break
                    else:
                        got_ranked, gt_item, alt_id = compare_with_multiple_gt_with_altid(self.kb, candidate, gt_list)
                        if got_ranked:
                            rank = idx
                            break

            else:
                ground_truth_id = item["ground_truth"]["id"]
                retrieved_candidates = item["retrieved_candidates"]
                rank = 0
                for idx, candidate in enumerate(retrieved_candidates, start=1):
                    if candidate["id"] == ground_truth_id:
                        rank = idx
                        break
            
            # Calculate reciprocal rank; if not found, reciprocal rank is 0
            reciprocal_rank = 1 / rank if rank > 0 else 0
            reciprocal_ranks.append(reciprocal_rank)
        
        # Calculate MRR by averaging all reciprocal ranks
        mrr = sum(reciprocal_ranks) / len(reciprocal_ranks)
        self.text_report+=f"\n\nCalculated MRR for {len(data)} items. \nSo, MRR : {mrr}\n\n"
        return mrr

    def calculate_recall_at_k(self, k):
        data = self.not_none_data
        relevant_count = 0
        for item in data:
            ground_truth_id = item["ground_truth"]["id"]
            retrieved_candidates = item["retrieved_candidates"][:k]  # Consider top-k candidates
            
            # Check if the relevant item is within the top-k candidates
            if any(candidate["id"] == ground_truth_id for candidate in retrieved_candidates):
                relevant_count += 1
        
        # Calculate recall@k
        recall_at_k = relevant_count / len(data)
        self.text_report+=f"\n\nFound {relevant_count} items out of {len(data)} within top-{k} candidates \nSo, recall@{k} : {relevant_count}/{len(data)}={recall_at_k}\n\n"
        return recall_at_k
    

    def category_specific_count(self, item, matched_gt, gt_list, is_matched, k):
        sample_id=item["sample_id"]
        mention=item["mention"]
        candidates = ' | '.join([i['title'] for i in item[self.candidate_key]])

        for gti in  gt_list:
            if gti['id'] in self.kb_id_list:
                title = gti['title']
                break
            else:
                title = gti['title']




        cat = get_category(mention, title)
        self.cat_wise_gt_matched[k][cat]['count'] += 1
        self.cat_wise_gt_matched[k][cat]['items'].append({
                'mention_id':sample_id, 
                'mention':mention, 
                'gt_title':title, 
                'matched': is_matched,
                'retrieved_candidates': candidates
                })
        if is_matched:
                self.cat_wise_gt_matched[k][cat]['matched'] += 1

    def calculate_recall_at_k_for_graph_retriever(self, k):
        data = self.not_none_data
        self.gt_matched = []
        self.gt_did_not_matched = []

        self.cat_wise_gt_matched[k] = {
            "HO":{'count':0, 'matched':0, 'items':[]}, 
            "MINT":{'count':0, 'matched':0, 'items':[]}, 
            "LO":{'count':0, 'matched':0, 'items':[]},
            "NO":{'count':0, 'matched':0, 'items':[]}
        }
        
        self.mention_matched_at_k[k] = []
        
        matched_count = 0

        if k==10:
            self.mention_wise_triple_matched[k]=[]

        for item in data:
            gt_list = self.multiple_gt_grag_dict[item['sample_id']]['ground_truth']
            unique_candidates = {}
            for cand in item[self.candidate_key]:
                unique_candidates[cand['id']] = cand
            item["unique_candidates"] = unique_candidates

            retrieved_candidates = list(item["unique_candidates"].keys())
            top_k = retrieved_candidates[:k] # Consider top-k candidates
            
            is_matched = False
            matched_gt = {}
            # Check if the relevant item is within the top-k candidates
            for i, candidate in enumerate(top_k):
                candidate_info = item["unique_candidates"][candidate]
                got_match_for_multi_gt, matched_gt = compare_with_multiple_gt(candidate_info, gt_list)
                if got_match_for_multi_gt:
                    matched_count += 1
                    is_matched = True
                    self.mention_matched_at_k[k].append(item["sample_id"])
                    self.gt_matched.append({
                        'sample_id':item["sample_id"],
                        'mention':item["mention"],
                        'mention_context':item["mention_context"],
                        'ground_truth':gt_list,
                        'rank':i+1,
                        'matched_triple': {**{'triple':candidate}, **candidate_info},
                        'unique_candidates':item["unique_candidates"]
                        })
                    break
                else:
                    got_match_for_multi_altid, matched_gt, altid = compare_with_multiple_gt_with_altid(self.kb, candidate_info, gt_list)
                    if got_match_for_multi_altid:
                        matched_count += 1
                        is_matched = True
                        self.mention_matched_at_k[k].append(item["sample_id"])
                        self.gt_matched.append({
                            'sample_id':item["sample_id"],
                            'mention':item["mention"],
                            'mention_context':item["mention_context"],
                            'ground_truth':gt_list,
                            'rank':i+1,
                            'matched_triple': {**{'triple':candidate}, **candidate_info},
                            'unique_candidates':item["unique_candidates"]
                            })
                        break

            self.category_specific_count(item, matched_gt, gt_list, is_matched=is_matched, k=k)
            
            if not is_matched:
                self.gt_did_not_matched.append({
                            'sample_id':item["sample_id"],
                            'mention':item["mention"],
                            'mention_context':item["mention_context"],
                            'ground_truth':gt_list,
                            'rank':-1,
                            'matched_triple':None,
                            'unique_candidates':item["unique_candidates"]
                            })
        
        # Calculate recall@k
        recall_at_k = matched_count / len(data)
        self.text_report+=f"\n\nFound {matched_count} items out of {len(data)} within top-{k} candidates \nSo, recall@{k} : {matched_count}/{len(data)}={recall_at_k}\n\n"
        # self.text_report+=f"\n\n"
        for cat in self.cat_wise_gt_matched[k]:
            count = self.cat_wise_gt_matched[k][cat]['count']
            matched = self.cat_wise_gt_matched[k][cat]['matched']
            score = matched / count if count > 0 else 0
            self.text_report+=f"{cat}, count: {count}, matched: {matched}, score: {score}\n"

        return recall_at_k


def samples_test_ent_appears_in_train(train_file, test_file):
    with open(train_file, 'r') as f:
        train = json.load(f)
    train_ents = {}
    for td in train:
        train_ents[td['mention_data']['kb_id']] = td

    count_appers = 0
    count_didnt_appers = 0
    unseen = []
    with open(test_file, 'r') as f:
        test = json.load(f)
    for td in test:
        test_id = td['mention_data']['kb_id']
        if test_id in train_ents:
            count_appers+=1
        else:
            count_didnt_appers+=1
            unseen.append(test_id)

    print(f'train : {len(train)}\ntest : {len(test)}\n{count_appers} GT entities from test set is also appears as GT in train set')
    return count_didnt_appers, unseen




def cat_eval(args, acc, filepath, kbpath):
    with open(filepath) as f:
        preds = json.load(f)
    with open(kbpath) as f:
        kb = json.load(f)
    
    count_didnt_appers, unseen = samples_test_ent_appears_in_train(
        args.train_data, args.lego_test_data)
    
    converted = []
    converted_unseen = []
    for p in preds:
        mention_context = p['text']
        gtid = p['mention_data']['kb_id']
        gttitle = kb[gtid]['title']
        retrieved_candidates = p['retrieved_candidates']
        reranker_result_gt = {}

        for i, c in enumerate(retrieved_candidates):
            if c['id'] == gtid:
                reranker_result_gt = {
                    "gt": c['id'],
                    "score": c['score'],
                    "rank": i+1
                }
                break

        candidates = [{'id':i, 'title':''} for i in p['mention_data']['candidates'] if i!=gtid ]
    
        mention = mention_context.split("[E1]")[1].split("[\E1]")[0].strip()
        
        d = {'mention_id' : '111111',
            'sample_id':p['sample_id'],
            'mention': mention,
            'mention_context':mention_context,
            'ground_truth': {'id':gtid, 'title':gttitle},
            'reranker_result_gt':reranker_result_gt,
            'retrieved_candidates':retrieved_candidates,
            'blink_retrieved_candidates':candidates
            }
        
        converted.append(d)

        if gtid in unseen:
            converted_unseen.append(d)


    with open(f'{args.grag_data_path}/test_grag.json') as f:
        multiple_gt_grag = json.load(f)
        multiple_gt_grag_dict = {}
        for item in multiple_gt_grag:
            multiple_gt_grag_dict[item['sample_id']] = item

    with open(f'{args.kb_file_path}') as f:
        exact_kb = json.load(f)

    eval_all = MultiGTEvaluation(converted, exact_kb, 'retrieved_candidates', multiple_gt_grag, for_retrieval=True)
    not_none_data, none_data = eval_all.get_report()
    eval_unseen = MultiGTEvaluation(converted, exact_kb, 'retrieved_candidates', multiple_gt_grag, for_retrieval=False)
    not_none_data, none_data = eval_unseen.get_report()

    # eval_all = Evaluation(converted, for_retrieval=False)
    # not_none_data, none_data = eval_all.get_report()
    # eval_unseen = Evaluation(converted_unseen, for_retrieval=False)
    # not_none_data, none_data = eval_unseen.get_report()

    report = f'Overall Accuracy : {acc}\n\n{eval_all.text_report}'
    report += f'\n\n\n\nAccuracy for {count_didnt_appers} unseen test sample'
    report += f'\n{eval_unseen.text_report}'

    eval_file_name = f"{filepath.replace('.json', '_eval.txt')}"
    eval_category_info = f"{filepath.replace('.json', '_category_info.json')}"
    eval_grag = f"{filepath.replace('.json', '_grag.json')}"

    if args.eval_before_fine_tune:
        eval_file_name = f"{filepath.replace('.json', '_eval_before_fine_tune.txt')}"
        eval_category_info = f"{filepath.replace('.json', '_category_info_before_fine_tune.json')}"
        eval_grag = f"{filepath.replace('.json', '_grag_before_fine_tune.json')}"
        
    with open(eval_file_name, 'w') as f:
        f.write(report)
    with open(eval_category_info, 'w') as f:
        json.dump(eval_all.cat_wise_gt_matched, f, indent=1)

    with open(eval_grag, 'w') as f:
        json.dump(converted, f, indent=1)
    



def eval_res_dataset(filepath, kbpath):
    with open(filepath) as f:
        preds = json.load(f)
    with open(kbpath) as f:
        kb = json.load(f)
    converted = []
    for p in preds:
        mention_context = p['text']
        gtid = p['mention_data']['kb_id']
        gttitle = kb[gtid]['title']
        candidates = [{'id':i, 'title':kb[i]['title']} for i in p['mention_data']['candidates']][:63]
        unique_triple = {}
        for c in candidates:
            unique_triple[c['id']]=c

        mention = mention_context.split("[E1]")[1].split("[\E1]")[0].strip()
        d = {'mention_id' : '111111',
            'mention': mention,
            'mention_context':mention_context,
            'ground_truth': {'id':gtid, 'title':gttitle},
            'retrieved_candidates':candidates,
            'unique_triple':unique_triple}
        converted.append(d)

    with open(f"{filepath.replace('.json', '_grag.json')}", 'w') as f:
        json.dump(converted, f, indent=1)
        
    eval_all = Evaluation(converted, for_retrieval=True)
    not_none_data, none_data = eval_all.get_report()
    report = f'{eval_all.text_report}'
    with open(f"{filepath.replace('.json', '_eval.txt')}", 'w') as f:
        f.write(report)
    
    with open(f"{filepath.replace('.json', '_category_info.json')}", 'w') as f:
        json.dump(eval_all.cat_wise_gt_matched, f, indent=1)


class PPTReport:
    def __init__(self):
        self.presentation = Presentation()

    def save_pptx(self, file_path):
        self.presentation.save(f"{file_path}.pptx")

    
    def add_new_slide(self, df, title):
        slide = self.presentation.slides.add_slide(self.presentation.slide_layouts[5])
        slide.shapes.title.text = title
        x, y, cx, cy = Inches(1), Inches(1.6), Inches(8), Inches(2)
        rows, cols = df.shape[0] + 1, df.shape[1]
        table = slide.shapes.add_table(rows, cols, x, y, cx, cy).table

        for table_col in table.columns:
            table_col.width = Inches(1)

        table.columns[0].width = Inches(2)

        for col_idx, column_name in enumerate(df.columns):
            table.cell(0, col_idx).text = column_name

        for row_idx, row in df.iterrows():
            for col_idx, value in enumerate(row):
                table.cell(row_idx + 1, col_idx).text = str(value)

        
class ReportMaker():
    def __init__(self, onto, report_save_to, report_start_end_text,
                firt_row_text, train_data_path, before_ft_dir=None): 
        self.onto = onto
        self.report_save_to = report_save_to
        self.report_start_end_text_for_recall = report_start_end_text
        self.firt_row_text = firt_row_text
        self.before_ft_dir = before_ft_dir
        self.train_data_path = train_data_path


    def get_cat_wise_score(self, report, row):
        lines = report.split("\n")
        for line in lines:
            if "HO," in line:
                l = line.split('score:')
                lr = l[1].strip()
                row['HO'] =  round(float(lr) * 100, 2)
                ho_count = int(re.search(r'count:\s*(\d+)', line).group(1))
                ho_matched = int(re.search(r'matched:\s*(\d+)', line).group(1))
            elif "MINT," in line:
                l = line.split('score:')
                lr = l[1].strip()
                row['MINT'] = round(float(lr) * 100, 2)
                mint_count = int(re.search(r'count:\s*(\d+)', line).group(1))
                mint_matched = int(re.search(r'matched:\s*(\d+)', line).group(1))
            elif "LO," in line:
                l = line.split('score:')
                lr = l[1].strip()
                row['LO'] = round(float(lr) * 100, 2)
            elif "NO," in line:
                l = line.split('score:')
                lr = l[1].strip()
                row['NO'] = round(float(lr) * 100, 2)

        # merge HO & MINT
        ho_mint_count = ho_count+mint_count
        ho_mint_matched = ho_matched+mint_matched

        row['HO'] = round((ho_mint_matched/ho_mint_count) * 100, 2)
        del row['MINT']


    # def get_cat_wise_score(self, report, row):
    #     lines = report.split("\n")
    #     for line in lines:
    #         if "HO," in line:
    #             l = line.split('score:')
    #             lr = l[1].strip()
    #             row['HO'] =  round(float(lr) * 100, 2)
    #         elif "MINT," in line:
    #             l = line.split('score:')
    #             lr = l[1].strip()
    #             row['MINT'] = round(float(lr) * 100, 2)
    #         elif "LO," in line:
    #             l = line.split('score:')
    #             lr = l[1].strip()
    #             row['LO'] = round(float(lr) * 100, 2)
    #         elif "NO," in line:
    #             l = line.split('score:')
    #             lr = l[1].strip()
    #             row['NO'] = round(float(lr) * 100, 2)


    def save_to_wandb(self, log_data):
        for entry in log_data:
            wandb.log({
                "Overall": entry["Overall"],
                "HO": entry["HO"],
                # "MINT": entry["MINT"],
                "LO": entry["LO"],
                "NO": entry["NO"],
                'MRR':entry["MRR"],
            }, step=entry["Epoch"])

    def get_mrr_and_overall(self, epoch, report, best_epoch):
        lines = report.split("\n")
        for line in lines:
            if "MRR :" in line:
                l = line.split('MRR :')
                lr = l[1].strip()
                mrr =  round(float(lr) * 100, 2)

            if "recall@1" in line:
                l = line.split('recall@1')
                lr = l[1].split('=')
                overall_acc =  round(float(lr[1]) * 100, 2)
                if overall_acc>best_epoch['acc']:
                    best_epoch = {'acc':overall_acc, 'report':report, 'epoch':epoch, 'MRR':mrr}

         
        return mrr, overall_acc, best_epoch
         


            

    # def get_best_epoch_and_send_all_wandb(self, exp, eval_text_name, text_eval_file_dir, epoch_start, epoch_end):
        
    #     report_all = []
    #     if self.before_ft_dir:
    #         row = self.before_ft_results
    #         report_all.append(row)

    #     best_epoch = {'acc':0, 'report':'', 'epoch':None, 'MRR':None}
    #     for i in range(epoch_start, epoch_end+1):
    #         epoch_dir = f'{text_eval_file_dir}{exp}/{self.onto}_{exp}_{i}pred/'
            
    #         try:
    #             with open(f"{epoch_dir}{eval_text_name}", 'r') as f:
    #                 report = f.read()
    #                 report = report.split(self.report_start_end_text[0])[1]
    #                 report = report.split(self.report_start_end_text[1])[0]
    #         except FileNotFoundError:
    #             print(f"No report found : \n{epoch_dir}")
    #             continue

    #         if epoch_start==0:
    #             epoch = i+1
    #         elif epoch_start==1:
    #             epoch = i

    #         mrr, overall_acc, best_epoch = self.get_mrr_and_overall(epoch, report, best_epoch)

    #         row = {'Epoch':epoch,'Overall': overall_acc, 'MRR':mrr}
    #         self.get_cat_wise_score(report, row)
    #         report_all.append(row)

    #     wandb.init(dir="wandb_logs", 
    #                project="fine-tune ReS", 
    #                name=exp)
    #     self.save_to_wandb(report_all)
    #     wandb.finish()

    #     return best_epoch



    def get_best_epoch_and_send_all_wandb(self, exp, eval_text_name, text_eval_file_dir, epoch_start, epoch_end):
        
        report_all = []
        if self.before_ft_dir:
            row = self.before_ft_results
            report_all.append(row)
        main_recall = 'recall@1'
        re_st_end = self.report_start_end_text_for_recall[main_recall]
        best_epoch = {'acc':0, 'report':'', 'epoch':None, 'MRR':None}
        entire_report_best_model = ''
        for i in range(epoch_start, epoch_end+1):
            # epoch_dir = f'{text_eval_file_dir}{exp}/epoch_{i}/'
            epoch_dir = f'{text_eval_file_dir}{exp}/{self.onto}_{exp}_{i}pred/'
            try:
                with open(f"{epoch_dir}{eval_text_name}", 'r') as f:
                    entire_report = f.read()
                    report = entire_report.split(re_st_end[0])[1]
                    report = report.split(re_st_end[1])[0]
            except FileNotFoundError:
                print(f"No report found : \n{epoch_dir}")
                return None

            # if epoch_start==0:
            #     epoch = i+1
            # elif epoch_start==1:

            epoch = i

            mrr, overall_acc, best_epoch, was_best = self.get_mrr_and_overall_for_best_model(epoch, report, best_epoch, 
                                                                                   main_recall)
            if was_best:
                entire_report_best_model = entire_report
                                                                                   

            row = {'Epoch':epoch,'Overall': overall_acc, 'MRR':mrr}
            self.get_cat_wise_score(report, row)
            report_all.append(row)

        # wandb.init(dir="wandb_logs", 
        #            project=f"FT-reranker-BLINK-{self.onto}", 
        #            name=exp+'20neg')
        

        # self.save_to_wandb(report_all)
        # wandb.finish()

        self.get_other_reacalls(entire_report_best_model, best_epoch)

        return best_epoch
    
    def get_specific_epoch_for_specifc_exp(self, exp, eval_text_name, text_eval_file_dir, epoch_start, epoch_end):

        exp_epoch_dict = {
            "(m1_e1)U(m3_e1)_multi_primeU(m4_e2)_multi_prime":1,
            "(m1_e1)U(m3_e1)_multi_prime_rm_sm_eU(m4_e2)_multi_prime_rm_sm_e":2,
            "synonym":4,
            "synonymU(m1_e1)U(m3_e1)_multi_prime_rm_sm_eU(m4_e2)_multi_prime_rm_sm_e":4
        }
        epoch_start, epoch_end = exp_epoch_dict[exp], exp_epoch_dict[exp]
        print(f'Epoch : {exp_epoch_dict[exp]} : {exp}')
        
        report_all = []
        if self.before_ft_dir:
            row = self.before_ft_results
            report_all.append(row)
        main_recall = 'recall@1'
        re_st_end = self.report_start_end_text_for_recall[main_recall]
        best_epoch = {'acc':0, 'report':'', 'epoch':None, 'MRR':None}
        entire_report_best_model = ''
        for i in range(epoch_start, epoch_end+1):
            # epoch_dir = f'{text_eval_file_dir}{exp}/epoch_{i}/'
            epoch_dir = f'{text_eval_file_dir}{exp}/{self.onto}_{exp}_{i}pred/'
            try:
                with open(f"{epoch_dir}{eval_text_name}", 'r') as f:
                    entire_report = f.read()
                    report = entire_report.split(re_st_end[0])[1]
                    report = report.split(re_st_end[1])[0]
            except FileNotFoundError:
                print(f"No report found : \n{epoch_dir}")
                return None

            epoch = i

            mrr, overall_acc, best_epoch, was_best = self.get_mrr_and_overall_for_best_model(epoch, report, best_epoch, 
                                                                                   main_recall)
            if was_best:
                entire_report_best_model = entire_report
                                                                                   

            row = {'Epoch':epoch,'Overall': overall_acc, 'MRR':mrr}
            self.get_cat_wise_score(report, row)
            report_all.append(row)

        # wandb.init(dir="wandb_logs", 
        #            project=f"FT-reranker-BLINK-{self.onto}", 
        #            name=exp+'20neg')
        

        # self.save_to_wandb(report_all)
        # wandb.finish()

        self.get_other_reacalls(entire_report_best_model, best_epoch)

        return best_epoch
    
    def get_other_reacalls(self, entire_report_best_model, best_epoch):
        recalls = {}
        for reacall_at_i in self.report_start_end_text_for_recall:
            if reacall_at_i == 'recall@1':
                continue
            report_start_end_text = self.report_start_end_text_for_recall[reacall_at_i]
            recall_report = entire_report_best_model.split(report_start_end_text[0])[1]
            recall_report = recall_report.split(report_start_end_text[1])[1]
            recall_report = recall_report.split(report_start_end_text[2])[0]
            recall_row = {'Overall': self.get_overall(recall_report, reacall_at_i)}
            self.get_cat_wise_score(recall_report, recall_row)
            recalls[reacall_at_i] = recall_row
        best_epoch['recalls'] = recalls

    # def get_report_before_fine_tune_result(self, eval_text_name):
    #     try:
    #         with open(f"{self.before_ft_dir}{eval_text_name}", 'r') as f:
    #             report = f.read()
    #             report = report.split(self.report_start_end_text[0])[1]
    #             report = report.split(self.report_start_end_text[1])[0]
    #     except FileNotFoundError:
    #         print(f"No report found : \n{self.before_ft_dir}")
    #     mrr, overall_acc, best_epoch = self.get_mrr_and_overall(0, report,  
    #     {'acc':0, 'report':'', 'epoch':None, 'MRR':None})
    #     row = {'Epoch':0,'Overall': overall_acc, 'MRR':mrr}
    #     self.get_cat_wise_score(report, row)
    #     return row

    def get_mrr_and_overall_for_best_model(self, epoch, report, best_epoch, reacall_at_i):
        mrr = None
        was_best = False
        lines = report.split("\n")
        for line in lines:
            if "MRR :" in line:
                l = line.split('MRR :')
                lr = l[1].strip()
                mrr =  round(float(lr) * 100, 2)
            
            if reacall_at_i in line:
                l = line.split(reacall_at_i)
                lr = l[1].split('=')
                overall_acc =  round(float(lr[1]) * 100, 2)
                if overall_acc>best_epoch['acc']:
                    best_epoch = {'acc':overall_acc, 'report':report, 'epoch':epoch, 'MRR':mrr}
                    was_best = True
        return mrr, overall_acc, best_epoch, was_best
         
    def get_overall(self, report, reacall_at_i):
        lines = report.split("\n")
        for line in lines:
            if reacall_at_i in line:
                l = line.split(reacall_at_i)
                lr = l[1].split('=')
                overall_acc =  round(float(lr[1]) * 100, 2)
                return overall_acc
            
    def get_report_before_fine_tune_result(self, eval_text_name):

        main_recall = 'recall@1'
        re_st_end = self.report_start_end_text_for_recall[main_recall]
        try:
            with open(f"{self.before_ft_dir}{eval_text_name}", 'r') as f:
                entire_report = f.read()
                report = entire_report.split(re_st_end[0])[1]
                report = report.split(re_st_end[1])[0]
        except FileNotFoundError:
            print(f"No report found : \n{self.before_ft_dir}")
            
        mrr, overall_acc, best_epoch, was_best = self.get_mrr_and_overall_for_best_model(0, report,  
        {'acc':0, 'report':'', 'epoch':None, 'MRR':None}, main_recall)
        row = {'Epoch':0,'Overall': overall_acc, 'MRR':mrr}
        self.get_cat_wise_score(report, row)

        recalls = {}
        for reacall_at_i in self.report_start_end_text_for_recall:
            if reacall_at_i == 'recall@1':
                continue
            report_start_end_text = self.report_start_end_text_for_recall[reacall_at_i]

            recall_report = entire_report.split(report_start_end_text[0])[1]
            recall_report = recall_report.split(report_start_end_text[1])[1]
            recall_report = recall_report.split(report_start_end_text[2])[0]
            recall_row = {'Overall': self.get_overall(recall_report, reacall_at_i)}
            self.get_cat_wise_score(recall_report, recall_row)
            recalls[reacall_at_i] = recall_row

        row['recalls'] = recalls
        
        return row

    # def make_pptx_report_for_all_epoch(self,text_eval_file_dir, eval_text_name, exps, epoch_start, epoch_end):
        
    #     ppt_report = PPTReport()
    #     reports = []

    #     if self.before_ft_dir:
    #         self.before_ft_results = self.get_report_before_fine_tune_result(eval_text_name)
    #         row = {**{f'Settings : {self.firt_row_text}':f" Original model (without fine-tune)​"},
    #             **self.before_ft_results }
    #         del row['Epoch']
    #         reports.append(row)


    #     for exp in exps:
    #         best_epoch = self.get_best_epoch_and_send_all_wandb(exp, eval_text_name, 
    #                         text_eval_file_dir, epoch_start, epoch_end)
            
    #         with open(f'{self.train_data_path}{exp}/train.json') as f:
    #             num_samples = len(json.load(f))

    #         row = {f'Settings : {self.firt_row_text}':f"{exp},({num_samples}),epoch {best_epoch['epoch']}", 
    #                'Overall':best_epoch['acc'], 'MRR':best_epoch['MRR']}
    #         self.get_cat_wise_score(best_epoch['report'], row)
    #         reports.append(row)

    #     df = pd.DataFrame(reports)
    #     title = f"{self.onto.upper()}:  Without fine-tuning retriever, 64 candidates"
    #     ppt_report.add_new_slide(df, title)
    #     ppt_file_path = f"{self.report_save_to}/evaluation_report"
    #     ppt_report.save_pptx(ppt_file_path)


    # def make_pptx_report_for_all_epoch(self,text_eval_file_dir, eval_text_name, exps, epoch_start, epoch_end,
    #                                    exp_from=None):

        
        
    #     ppt_report = PPTReport()
        

       
    #     recall_at_1_reports = []
        
    #     recall_at_n_reports = {}
    #     for reacall_at_i in self.report_start_end_text_for_recall:
    #         recall_at_n_reports[reacall_at_i]=[]

    #     if self.before_ft_dir:
    #         self.before_ft_results = self.get_report_before_fine_tune_result(eval_text_name)
    #         row = {**{f'Settings : {self.firt_row_text}':f" Original model (without fine-tune)​"},
    #             **self.before_ft_results }
    #         del row['Epoch']
            

    #         for recall_i in row['recalls']:
    #             recall_i_data = {"Settings : ":"Original model (without fine-tune)​"}
    #             for key in row['recalls'][recall_i]:
    #                 recall_i_data[key] = row['recalls'][recall_i][key]
    #             recall_at_n_reports[recall_i].append(recall_i_data)


    #         del row['recalls']
    #         recall_at_1_reports.append(row)

    #     self.best_epoch_dir = {}
    #     for exp in exps:
    #         best_epoch = self.get_best_epoch_and_send_all_wandb(exp, eval_text_name, 
    #                         text_eval_file_dir, epoch_start, epoch_end)
            

    #         if epoch_start==0:
    #             best_epoch_num = best_epoch['epoch']-1
    #         elif epoch_start==1:
    #             best_epoch_num = best_epoch['epoch']

    #         self.best_epoch_dir[exp] = f"epoch_{best_epoch_num}"
            
    #         # if not exp_from:
    #         exp_from = exp

    #         # num_samples =len( read_jsonl(f'data/{self.onto}/blink_format/{self.sptitname}/{exp_from}/train.jsonl'))
    #         with open(f'{self.train_data_path}{exp}/train.json') as f:
    #             num_samples = len(json.load(f))
                
    #         # cat_stats = self.get_cat_stats(exp_from)

    #         setting_name = f'Settings : {self.firt_row_text}'
    #         row = {setting_name:f"{exp}, ({num_samples}), epoch {best_epoch['epoch']}", 
    #             'Overall':best_epoch['acc'], 'MRR':best_epoch['MRR']}
    #         self.get_cat_wise_score(best_epoch['report'], row)
            
            

    #         for recall_i in best_epoch['recalls']:
    #             recall_i_data = {"Settings : ":row[setting_name]}
    #             for key in best_epoch['recalls'][recall_i]:
    #                 recall_i_data[key] = best_epoch['recalls'][recall_i][key]
                
    #             recall_at_n_reports[recall_i].append(recall_i_data)

            
    #         recall_at_1_reports.append(row)


        
    #     df = pd.DataFrame(recall_at_1_reports)
    #     title = f"{self.onto.upper()}:  recall@1"
    #     ppt_report.add_new_slide(df, title)
    
    #     for reacall_at_i in self.report_start_end_text_for_recall:
    #         if reacall_at_i == 'recall@1':
    #             continue

    #         df = pd.DataFrame(recall_at_n_reports[reacall_at_i])
    #         title = f"{self.onto.upper()}:  {reacall_at_i}"
    #         ppt_report.add_new_slide(df, title)


    #     ppt_file_path = f"{self.report_save_to}/evaluation_report"
    #     ppt_report.save_pptx(ppt_file_path)



    def make_pptx_report_for_all_epoch(self, text_eval_file_dir, eval_text_name, exps,
                                    epoch_start, epoch_end, exp_from=None):
        """Build PPTX (one slide per recall@k) and a matching JSON dump.

        Note: recall@k>1 cat-wise values come from
        self.get_best_epoch_and_send_all_wandb(...). If those look wrong
        (e.g. HO/LO/NO pinned at 100.0 while Overall varies), the bug is in
        that upstream call — fix it where best_epoch['recalls'][recall@k]
        is populated, not here.
        """
        
        ppt_report = PPTReport()
        recall_at_1_reports = []
        recall_at_n_reports = {k: [] for k in self.report_start_end_text_for_recall}

        settings_key = f'Settings : {self.firt_row_text}'

        # ---- No-fine-tune row ----
        if self.before_ft_dir:
            self.before_ft_results = self.get_report_before_fine_tune_result(eval_text_name)
            row = {**{settings_key: " Original model (without fine-tune)"},
                **self.before_ft_results}
            # No-fine-tune has no training epoch; render as None (NaN in PPTX).
            row['Epoch'] = None

            for recall_i in row['recalls']:
                recall_i_data = {"Settings : ": "Original model (without fine-tune)",
                                "Epoch": None}
                for key in row['recalls'][recall_i]:
                    recall_i_data[key] = row['recalls'][recall_i][key]
                recall_at_n_reports[recall_i].append(recall_i_data)

            del row['recalls']
            recall_at_1_reports.append(row)

        # ---- Per-experiment rows ----
        self.best_epoch_dir = {}
        for exp in exps:
            best_epoch = self.get_best_epoch_and_send_all_wandb(
                exp, eval_text_name, text_eval_file_dir, epoch_start, epoch_end)

            # best_epoch = self.get_specific_epoch_for_specifc_exp(
            #     exp, eval_text_name, text_eval_file_dir, epoch_start, epoch_end)

            # if epoch_start == 0:
            #     best_epoch_num = best_epoch['epoch'] - 1
            # elif epoch_start == 1:
            #     best_epoch_num = best_epoch['epoch']
            
            best_epoch_num = best_epoch['epoch']

            self.best_epoch_dir[exp] = f"epoch_{best_epoch_num}"

            exp_from = exp
            with open(f'{self.train_data_path}{exp}/train.json') as f:
                num_samples = len(json.load(f))

            long_setting = f"{exp}, ({num_samples}), epoch {best_epoch['epoch']}"

            row = {settings_key: long_setting,
                'Epoch':   best_epoch['epoch'],
                'Overall': best_epoch['acc'],
                'MRR':     best_epoch['MRR']}
            self.get_cat_wise_score(best_epoch['report'], row)   # adds HO/LO/NO

            for recall_i in best_epoch['recalls']:
                recall_i_data = {"Settings : ": long_setting,
                                "Epoch": best_epoch['epoch']}
                for key in best_epoch['recalls'][recall_i]:
                    recall_i_data[key] = best_epoch['recalls'][recall_i][key]
                recall_at_n_reports[recall_i].append(recall_i_data)

            recall_at_1_reports.append(row)

        # ---- PPTX (one slide per recall@k) ----
        df = pd.DataFrame(recall_at_1_reports)
        ppt_report.add_new_slide(df, f"{self.onto.upper()}:  recall@1")

        for recall_i in self.report_start_end_text_for_recall:
            if recall_i == 'recall@1':
                continue
            df = pd.DataFrame(recall_at_n_reports[recall_i])
            ppt_report.add_new_slide(df, f"{self.onto.upper()}:  {recall_i}")

        ppt_report.save_pptx(f"{self.report_save_to}/evaluation_report")

        # ---- JSON dump ----
        # Backfill recall@1: the per-recall loop above iterates row['recalls'],
        # which only contains k>1. recall@1 metrics live at the top level of
        # each row, so we project them into the same shape as recall@k entries.
        recall_at_n_reports['recall@1'] = [
            {
                "Settings : ": r[settings_key],
                "Epoch":   r.get("Epoch"),
                "Overall": r.get("Overall"),
                "MRR":     r.get("MRR"),
                "HO":      r.get("HO"),
                "LO":      r.get("LO"),
                "NO":      r.get("NO"),
            }
            for r in recall_at_1_reports
        ]

        json_path = f"{self.report_save_to}/evaluation_report.json"
        with open(json_path, 'w') as f:
            json.dump(recall_at_n_reports, f, indent=2)


